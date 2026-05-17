"""
engine.py — 图片转可编辑 PPTX 核心引擎 v2

改进点（相比 v1）:
1. OCR 换为 PaddleOCR（检测率更高）
2. 颜色提取用 K-Means 聚类 + 多层回退
3. 字号用二分查找 + Pillow 渲染精确计算
4. 取消文本框合并，每个检测框独立
5. 禁用自动换行，防止文字跨行
"""

import time
import os
import re
import gc
import sys
import threading
from pathlib import Path
from typing import List, Tuple, Optional
from dataclasses import dataclass

import cv2
import numpy as np
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE


def _load_local_env():
    """Load local .env values for CLI and direct engine usage."""
    env_path = Path(__file__).parent / ".env"
    if not env_path.exists():
        return

    for line in env_path.read_text().splitlines():
        line = line.strip()
        if not line or line.startswith("#") or "=" not in line:
            continue
        key, val = line.split("=", 1)
        os.environ.setdefault(key.strip(), val.strip().strip('"').strip("'"))


_load_local_env()


def _get_torch():
    """Import torch only for features that actually need it."""
    import torch
    return torch


def _env_truthy(name: str, default: bool = False) -> bool:
    value = os.environ.get(name)
    if value is None:
        return default
    return value.strip().lower() in {"1", "true", "yes", "on"}


def _processing_limits() -> Tuple[int, int]:
    max_mp = float(os.environ.get("DECKLENS_MAX_PROCESS_MP", "3.0"))
    max_side = int(os.environ.get("DECKLENS_MAX_PROCESS_SIDE", "2200"))
    return max(1, int(max_mp * 1_000_000)), max(256, max_side)


def _scale_for_limits(width: int, height: int) -> float:
    max_pixels, max_side = _processing_limits()
    scale = min(1.0, max_side / max(width, height), (max_pixels / max(1, width * height)) ** 0.5)
    return max(0.01, scale)


def constrain_image_for_processing(image_path: str, output_dir: Optional[str] = None) -> str:
    """
    Downscale oversized inputs before OCR/inpainting/segmentation.

    SAM/Paddle/LaMa memory use grows with pixel count, and automatic mask
    generation can allocate many full-size masks. Keeping local processing near
    slide/screenshot resolution prevents accidental 30GB+ jobs.
    """
    image_path = str(Path(image_path).resolve())
    with Image.open(image_path) as img:
        width, height = img.size
        scale = _scale_for_limits(width, height)
        if scale >= 0.999:
            return image_path

        new_size = (max(1, int(round(width * scale))), max(1, int(round(height * scale))))
        if output_dir is None:
            output_dir = str(Path(image_path).parent)
        out_path = str(Path(output_dir) / f"{Path(image_path).stem}_decklens_scaled.png")
        resized = img.convert("RGB")
        resized.thumbnail(new_size, Image.Resampling.LANCZOS)
        resized.save(out_path, "PNG", optimize=True)

    print(f"  [MEM] 输入过大，已降采样: {width}x{height} -> {new_size[0]}x{new_size[1]}", flush=True)
    return out_path


# ─── 数据结构 ───

@dataclass
class TextBlock:
    """一个检测到的文字块"""
    text: str
    x: int
    y: int
    w: int
    h: int
    confidence: float = 0.0
    color: Tuple[int, int, int] = (0, 0, 0)
    font_size_pt: float = 12.0
    bold: bool = False
    font_name: str = "微软雅黑"
    original_x: Optional[int] = None
    original_y: Optional[int] = None
    original_w: Optional[int] = None
    original_h: Optional[int] = None

    def original_box(self) -> Tuple[int, int, int, int]:
        return (
            self.x if self.original_x is None else self.original_x,
            self.y if self.original_y is None else self.original_y,
            self.w if self.original_w is None else self.original_w,
            self.h if self.original_h is None else self.original_h,
        )


# ─── OCR 检测（PaddleOCR） ───

_paddle_ocr = None


def get_paddle_ocr():
    """懒加载 PaddleOCR"""
    global _paddle_ocr
    if _paddle_ocr is None:
        import warnings
        os.environ.setdefault("GLOG_minloglevel", "2")
        os.environ.setdefault("FLAGS_minloglevel", "2")
        warnings.filterwarnings("ignore")
        from paddleocr import PaddleOCR
        _paddle_ocr = PaddleOCR(
            lang="ch",
            use_doc_orientation_classify=False,
            use_doc_unwarping=False,
            use_textline_orientation=False,
        )
    return _paddle_ocr


def detect_text_paddle(image_path: str, expand_px: int = 3) -> List[TextBlock]:
    """
    用 PaddleOCR 检测文字，返回 TextBlock 列表。
    """
    ocr = get_paddle_ocr()
    image = cv2.imread(image_path)
    if image is None:
        raise ValueError(f"无法读取图片: {image_path}")

    img_h, img_w = image.shape[:2]
    result = ocr.predict(image_path)

    blocks = []
    if not result or len(result) == 0:
        return blocks

    page = result[0]
    dt_polys = page.get("dt_polys", [])
    rec_texts = page.get("rec_texts", [])
    rec_scores = page.get("rec_scores", [])

    for idx, (poly, text) in enumerate(zip(dt_polys, rec_texts)):
        score = rec_scores[idx] if idx < len(rec_scores) else 1.0
        if score < 0.5:
            continue
        if not text.strip():
            continue

        points = np.array(poly)
        x_min = max(0, int(np.min(points[:, 0])) - expand_px)
        y_min = max(0, int(np.min(points[:, 1])) - expand_px)
        x_max = min(img_w, int(np.max(points[:, 0])) + expand_px)
        y_max = min(img_h, int(np.max(points[:, 1])) + expand_px)

        w = x_max - x_min
        h = y_max - y_min
        if w <= 0 or h <= 0:
            continue

        # 过滤：面积太小或文字太短且不像正常文本的区域（可能是图标/装饰误识别）
        # 纯符号或单个非中文非数字字符，且面积小，跳过
        is_likely_noise = (
            len(text.strip()) <= 2
            and not re.search(r'[\u4e00-\u9fff0-9]', text)
            and w * h < img_w * img_h * 0.005
        )
        if is_likely_noise:
            continue

        # 提取颜色
        color = extract_text_color(image, x_min, y_min, x_max, y_max)

        # 计算字号（二分查找，只约束宽度）
        font_size = fit_font_size(text, w, h)

        # 粗体判断：大标题（字号 > 50pt）或字符平均宽度占比高
        # 参考 同类图片转 PPTX 工具：标题和强调文本标粗体
        avg_char_w = w / max(len(text), 1)
        is_large_title = font_size > 35
        is_wide_chars = avg_char_w > h * 0.75
        bold = is_large_title or is_wide_chars

        blocks.append(TextBlock(
            text=text,
            x=x_min, y=y_min, w=w, h=h,
            confidence=score,
            color=color,
            font_size_pt=font_size,
            bold=bold,
            original_x=x_min, original_y=y_min, original_w=w, original_h=h,
        ))

    return blocks


# ─── 字体匹配 ───

_CJK_RE = re.compile(r'[\u3400-\u4dbf\u4e00-\u9fff\uf900-\ufaff]')
_LATIN_RE = re.compile(r'[A-Za-z]')
_DIGIT_RE = re.compile(r'\d')

DEFAULT_CATEGORY_FONTS = {
    "zh": "微软雅黑",
    "en": "Arial",
    "num": "Arial",
}


def text_font_category(text: str) -> Optional[str]:
    """
    Classify text for page-level font unification.

    Mixed Chinese/English/digit text is treated as Chinese, because a single
    PowerPoint run needs a primary font and CJK glyph coverage is the safer base.
    """
    if _CJK_RE.search(text):
        return "zh"
    if _LATIN_RE.search(text):
        return "en"
    if _DIGIT_RE.search(text):
        return "num"
    return None


def normalize_page_fonts(blocks: List[TextBlock]) -> List[TextBlock]:
    """
    Keep at most three fonts per page: Chinese, English, and numeric.

    Each category uses the most common matched font on the page. Ties prefer the
    font covering the larger total text area, which makes titles and prominent
    labels influence the representative font without letting one-off matches
    create extra fonts.
    """
    stats = {}
    for block in blocks:
        category = text_font_category(block.text)
        if not category:
            continue
        font_name = block.font_name or DEFAULT_CATEGORY_FONTS[category]
        area = max(1, int(block.w) * int(block.h))
        category_stats = stats.setdefault(category, {})
        count, total_area = category_stats.get(font_name, (0, 0))
        category_stats[font_name] = (count + 1, total_area + area)

    selected = {}
    for category, category_stats in stats.items():
        if category_stats:
            selected[category] = max(category_stats.items(), key=lambda item: item[1])[0]
        else:
            selected[category] = DEFAULT_CATEGORY_FONTS[category]

    for block in blocks:
        category = text_font_category(block.text)
        if category:
            block.font_name = selected.get(category, DEFAULT_CATEGORY_FONTS[category])

    return blocks


def _rgb_to_lab(color: Tuple[int, int, int]) -> np.ndarray:
    rgb = np.array([[[
        max(0, min(255, int(color[0]))),
        max(0, min(255, int(color[1]))),
        max(0, min(255, int(color[2]))),
    ]]], dtype=np.uint8)
    return cv2.cvtColor(rgb, cv2.COLOR_RGB2LAB).astype(np.float32)[0, 0]


def _lab_color_distance(a: Tuple[int, int, int], b: Tuple[int, int, int]) -> float:
    return float(np.linalg.norm(_rgb_to_lab(a) - _rgb_to_lab(b)))


def _snap_font_size(size_pt: float) -> float:
    common_sizes = [6, 7, 8, 9, 10, 11, 12, 14, 16, 18, 20, 24, 28, 32, 36, 40, 44, 48, 54, 60, 72]
    nearest = min(common_sizes, key=lambda value: abs(value - size_pt))
    if abs(nearest - size_pt) <= max(0.75, size_pt * 0.06):
        return float(nearest)
    return float(round(size_pt))


def _font_size_close(a: TextBlock, b: TextBlock, tolerance: float = 0.16) -> bool:
    a_size = max(1.0, float(a.font_size_pt))
    b_size = max(1.0, float(b.font_size_pt))
    return abs(a_size - b_size) / max(a_size, b_size) <= tolerance


def _box_height_close(a: TextBlock, b: TextBlock, tolerance: float = 0.22) -> bool:
    a_h = max(1, int(a.h))
    b_h = max(1, int(b.h))
    return abs(a_h - b_h) / max(a_h, b_h) <= tolerance


def _style_layout_related(a: TextBlock, b: TextBlock, image_size: Optional[Tuple[int, int]] = None) -> bool:
    ax0, ay0, ax1, ay1 = a.x, a.y, a.x + a.w, a.y + a.h
    bx0, by0, bx1, by1 = b.x, b.y, b.x + b.w, b.y + b.h
    acx, acy = ax0 + a.w / 2, ay0 + a.h / 2
    bcx, bcy = bx0 + b.w / 2, by0 + b.h / 2
    avg_w = max(1.0, (a.w + b.w) / 2)
    avg_h = max(1.0, (a.h + b.h) / 2)

    x_tol = max(8.0, avg_w * 0.10)
    y_tol = max(6.0, avg_h * 0.45)
    left_aligned = abs(ax0 - bx0) <= x_tol
    center_aligned = abs(acx - bcx) <= x_tol
    right_aligned = abs(ax1 - bx1) <= x_tol
    row_aligned = abs(acy - bcy) <= y_tol
    column_related = left_aligned or center_aligned or right_aligned

    if row_aligned and abs(acx - bcx) <= max(avg_w * 5.0, 240.0):
        return True

    if column_related and abs(acy - bcy) <= max(avg_h * 8.0, 220.0):
        return True

    center_dist = ((acx - bcx) ** 2 + (acy - bcy) ** 2) ** 0.5
    if center_dist <= max(avg_h * 5.0, avg_w * 1.8, 140.0):
        return True

    if image_size:
        img_w, img_h = image_size
        same_horizontal_band = abs(acy - bcy) <= max(img_h * 0.035, y_tol)
        same_vertical_band = abs(acx - bcx) <= max(img_w * 0.035, x_tol)
        if same_horizontal_band or same_vertical_band:
            return True

    return False


def _style_layout_strongly_related(a: TextBlock, b: TextBlock) -> bool:
    ax0, ay0, ax1, ay1 = a.x, a.y, a.x + a.w, a.y + a.h
    bx0, by0, bx1, by1 = b.x, b.y, b.x + b.w, b.y + b.h
    acx, acy = ax0 + a.w / 2, ay0 + a.h / 2
    bcx, bcy = bx0 + b.w / 2, by0 + b.h / 2
    avg_w = max(1.0, (a.w + b.w) / 2)
    avg_h = max(1.0, (a.h + b.h) / 2)

    x_tol = max(10.0, avg_w * 0.12)
    y_tol = max(6.0, avg_h * 0.45)
    same_column = (
        abs(ax0 - bx0) <= x_tol
        or abs(acx - bcx) <= x_tol
        or abs(ax1 - bx1) <= x_tol
    )
    same_row = abs(acy - bcy) <= y_tol
    vertical_gap = max(0.0, max(ay0, by0) - min(ay1, by1))
    horizontal_gap = max(0.0, max(ax0, bx0) - min(ax1, bx1))

    if same_column and vertical_gap <= max(avg_h * 2.2, 80.0):
        return True
    if same_row and horizontal_gap <= max(avg_w * 2.0, 220.0):
        return True

    return False


def _cluster_by_value(items: List[TextBlock], value_fn, tolerance: float) -> List[List[TextBlock]]:
    ordered = sorted(items, key=value_fn)
    clusters: List[List[TextBlock]] = []
    for item in ordered:
        if not clusters:
            clusters.append([item])
            continue
        cluster_values = [value_fn(existing) for existing in clusters[-1]]
        center = float(np.median(cluster_values))
        if abs(value_fn(item) - center) <= tolerance:
            clusters[-1].append(item)
        else:
            clusters.append([item])
    return clusters


def _normalize_text_geometry_for_group(group_blocks: List[TextBlock]) -> None:
    """
    Stabilize geometry for repeated text that already shares one visual style.

    OCR boxes often jitter by a few pixels between items in the same menu/list.
    After font size normalization, those tiny bbox differences become visible in
    PPT because text is re-rendered inside independent text boxes.
    """
    if len(group_blocks) < 2:
        return

    heights = np.array([max(1, int(block.h)) for block in group_blocks], dtype=np.float32)
    median_h = max(1, int(round(float(np.median(heights)))))
    height_spread = float(np.max(np.abs(heights - median_h) / np.maximum(heights, median_h)))
    if height_spread <= 0.35:
        for block in group_blocks:
            center_y = block.y + block.h / 2
            block.h = median_h
            block.y = int(round(center_y - median_h / 2))

    x_tol = max(8.0, median_h * 0.35)
    y_tol = max(6.0, median_h * 0.35)

    snapped_x = set()
    for cluster in _cluster_by_value(group_blocks, lambda b: b.x, x_tol):
        if len(cluster) < 2:
            continue
        snapped_left = int(round(float(np.median([block.x for block in cluster]))))
        for block in cluster:
            block.x = snapped_left
            snapped_x.add(id(block))

    for cluster in _cluster_by_value(group_blocks, lambda b: b.x + b.w / 2, x_tol):
        cluster = [block for block in cluster if id(block) not in snapped_x]
        if len(cluster) < 2:
            continue
        snapped_center = float(np.median([block.x + block.w / 2 for block in cluster]))
        for block in cluster:
            block.x = int(round(snapped_center - block.w / 2))
            snapped_x.add(id(block))

    for cluster in _cluster_by_value(group_blocks, lambda b: b.x + b.w, x_tol):
        cluster = [block for block in cluster if id(block) not in snapped_x]
        if len(cluster) < 2:
            continue
        snapped_right = int(round(float(np.median([block.x + block.w for block in cluster]))))
        for block in cluster:
            block.x = int(round(snapped_right - block.w))

    for cluster in _cluster_by_value(group_blocks, lambda b: b.y + b.h / 2, y_tol):
        if len(cluster) < 2:
            continue
        snapped_center = float(np.median([block.y + block.h / 2 for block in cluster]))
        for block in cluster:
            block.y = int(round(snapped_center - block.h / 2))


def _block_intersection(a: TextBlock, b: TextBlock) -> Tuple[int, int, int, int]:
    return _box_intersection((a.x, a.y, a.w, a.h), (b.x, b.y, b.w, b.h))


def _box_intersection(
    a: Tuple[float, float, float, float],
    b: Tuple[float, float, float, float],
) -> Tuple[float, float, float, float]:
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    x0 = max(float(ax), float(bx))
    y0 = max(float(ay), float(by))
    x1 = min(float(ax + aw), float(bx + bw))
    y1 = min(float(ay + ah), float(by + bh))
    if x1 <= x0 or y1 <= y0:
        return (0, 0, 0, 0)
    return (x0, y0, x1 - x0, y1 - y0)


def _box_overlap_ratio(
    a: Tuple[float, float, float, float],
    b: Tuple[float, float, float, float],
) -> float:
    _, _, iw, ih = _box_intersection(a, b)
    if iw <= 0 or ih <= 0:
        return 0.0
    min_area = max(1.0, min(float(a[2] * a[3]), float(b[2] * b[3])))
    return float((iw * ih) / min_area)


def _fallback_text_width(text: str, font_size: float) -> float:
    width = 0.0
    for ch in text:
        if ch.isspace():
            width += font_size * 0.35
        elif _CJK_RE.search(ch):
            width += font_size
        elif ch.isdigit():
            width += font_size * 0.58
        elif ch.isupper():
            width += font_size * 0.66
        else:
            width += font_size * 0.56
    return max(font_size, width)


def _estimate_rendered_text_box(block: TextBlock) -> Tuple[float, float, float, float]:
    font_size = max(1.0, float(block.font_size_pt))
    text = block.text or ""
    text_w = _fallback_text_width(text, font_size)
    text_h = font_size * 1.18
    font_path = _get_font_path()
    if font_path:
        try:
            font = ImageFont.truetype(font_path, max(1, int(round(font_size))))
            bbox = font.getbbox(text)
            text_w = max(1.0, float(bbox[2] - bbox[0]) * 1.08)
            text_h = max(1.0, float(bbox[3] - bbox[1]) * 1.25)
        except Exception:
            pass

    visual_w = max(float(block.w), text_w)
    visual_h = max(float(block.h), text_h)
    visual_x = float(block.x)
    visual_y = float(block.y + block.h / 2 - visual_h / 2)
    return (visual_x, visual_y, visual_w, visual_h)


def prevent_small_text_overlaps(
    blocks: List[TextBlock],
    image_size: Optional[Tuple[int, int]] = None,
    generated_overlap_ratio: float = 0.03,
    original_large_overlap_ratio: float = 0.45,
    overlap_growth_ratio: float = 0.06,
) -> List[TextBlock]:
    """
    Remove accidental overlaps between final PPT text boxes.

    Each block keeps its original OCR bbox. We preserve pairs that were already
    heavily overlapped in the source image, then estimate the final rendered text
    footprint after font/style normalization and only move pairs that newly
    overlap in the generated PPT geometry.
    """
    if len(blocks) < 2:
        return blocks

    img_w, img_h = image_size or (0, 0)
    page_w = img_w or max(block.x + block.w for block in blocks)
    page_h = img_h or max(block.y + block.h for block in blocks)
    max_adjust_x = max(12.0, page_w * 0.12)
    max_adjust_y = max(8.0, page_h * 0.08)
    ordered = sorted(blocks, key=lambda b: (b.y + b.h / 2, b.x + b.w / 2))

    for _ in range(3):
        changed = False
        for i, a in enumerate(ordered):
            for b in ordered[i + 1:]:
                original_overlap = _box_overlap_ratio(a.original_box(), b.original_box())
                if original_overlap >= original_large_overlap_ratio:
                    continue

                a_visual = _estimate_rendered_text_box(a)
                b_visual = _estimate_rendered_text_box(b)
                ix, iy, iw, ih = _box_intersection(a_visual, b_visual)
                if iw <= 0 or ih <= 0:
                    continue

                generated_overlap = _box_overlap_ratio(a_visual, b_visual)
                if generated_overlap < generated_overlap_ratio:
                    continue
                if generated_overlap <= original_overlap + overlap_growth_ratio:
                    continue

                ax, ay, aw, ah = a_visual
                bx, by, bw, bh = b_visual
                acx, acy = ax + aw / 2, ay + ah / 2
                bcx, bcy = bx + bw / 2, by + bh / 2

                avg_h = max(1.0, (ah + bh) / 2)
                avg_w = max(1.0, (aw + bw) / 2)
                same_row = abs(acy - bcy) <= max(6.0, avg_h * 0.45)
                same_column = abs(acx - bcx) <= max(8.0, avg_w * 0.25)
                gap = max(1, int(round(min(avg_h, avg_w) * 0.08)))

                if same_row:
                    if bcx >= acx:
                        target = b
                        delta = int(round((ax + aw + gap) - bx))
                    else:
                        target = a
                        delta = int(round((bx + bw + gap) - ax))
                    if 0 < delta <= max_adjust_x:
                        target.x += delta
                        if img_w:
                            target.x = min(target.x, max(0, img_w - target.w))
                        changed = True
                elif same_column or ih <= max(2.0, min(ah, bh) * 0.35):
                    if bcy >= acy:
                        target = b
                        delta = int(round((ay + ah + gap) - by))
                    else:
                        target = a
                        delta = int(round((by + bh + gap) - ay))
                    if 0 < delta <= max_adjust_y:
                        target.y += delta
                        if img_h:
                            target.y = min(target.y, max(0, img_h - target.h))
                        changed = True

        if not changed:
            break
        ordered = sorted(blocks, key=lambda b: (b.y + b.h / 2, b.x + b.w / 2))

    return blocks


def normalize_text_styles(blocks: List[TextBlock], image_size: Optional[Tuple[int, int]] = None) -> List[TextBlock]:
    """
    Normalize repeated text styles on one page.

    OCR and per-block font matching are intentionally noisy on screenshots.
    This pass treats color, font size, and text category as hard constraints,
    then uses layout proximity/alignment to merge repeated blocks into a shared
    visual style.
    """
    if len(blocks) < 2:
        return blocks

    parent = list(range(len(blocks)))

    def find(i: int) -> int:
        while parent[i] != i:
            parent[i] = parent[parent[i]]
            i = parent[i]
        return i

    def union(a: int, b: int) -> None:
        ra = find(a)
        rb = find(b)
        if ra != rb:
            parent[rb] = ra

    for i, a in enumerate(blocks):
        for j in range(i + 1, len(blocks)):
            b = blocks[j]
            if text_font_category(a.text) != text_font_category(b.text):
                continue
            if _lab_color_distance(a.color, b.color) > 20.0:
                continue
            if not _box_height_close(a, b):
                continue
            layout_related = _style_layout_related(a, b, image_size=image_size)
            if not layout_related:
                continue
            if not _font_size_close(a, b) and not _style_layout_strongly_related(a, b):
                continue
            if a.bold != b.bold and max(a.font_size_pt, b.font_size_pt) >= 18:
                continue
            union(i, j)

    groups = {}
    for idx in range(len(blocks)):
        groups.setdefault(find(idx), []).append(idx)

    for indices in groups.values():
        if len(indices) < 2:
            continue

        group_blocks = [blocks[idx] for idx in indices]
        font_counts = {}
        for block in group_blocks:
            font_counts[block.font_name] = font_counts.get(block.font_name, 0) + 1
        font_name = max(font_counts.items(), key=lambda item: item[1])[0]

        size = _snap_font_size(float(np.median([block.font_size_pt for block in group_blocks])))
        color_values = np.array([block.color for block in group_blocks], dtype=np.float32)
        color = tuple(int(v) for v in np.median(color_values, axis=0))
        bold = sum(1 for block in group_blocks if block.bold) >= (len(group_blocks) / 2)

        for block in group_blocks:
            block.font_name = font_name
            block.font_size_pt = size
            block.color = color
            block.bold = bold
        _normalize_text_geometry_for_group(group_blocks)

    prevent_small_text_overlaps(blocks, image_size=image_size)
    return blocks


def match_fonts_for_blocks(image_path: str, blocks: List[TextBlock]) -> List[TextBlock]:
    """
    对每个文字块进行字体匹配，更新 block.font_name。
    从图片中裁剪出每个 block 的区域，调用 font_matcher 进行匹配。
    """
    import sys
    font_matcher_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "font_matcher")
    if font_matcher_dir not in sys.path:
        sys.path.insert(0, font_matcher_dir)

    from font_matcher import match_font

    image = cv2.imread(image_path)
    if image is None:
        return blocks

    img_h, img_w = image.shape[:2]

    for block in blocks:
        # 裁剪 block 区域
        x1 = max(0, block.x)
        y1 = max(0, block.y)
        x2 = min(img_w, block.x + block.w)
        y2 = min(img_h, block.y + block.h)

        if x2 <= x1 or y2 <= y1:
            continue

        region = image[y1:y2, x1:x2]
        # 转灰度
        if len(region.shape) == 3:
            region_gray = cv2.cvtColor(region, cv2.COLOR_BGR2GRAY)
        else:
            region_gray = region

        try:
            result = match_font(region_gray, block.text)
            block.font_name = result.get("pptx_name", "微软雅黑")
        except Exception:
            pass  # 匹配失败保持默认

    blocks = normalize_page_fonts(blocks)
    blocks = normalize_text_styles(blocks, image_size=(img_w, img_h))
    return blocks


# ─── 颜色提取（K-Means 聚类 + 多层回退） ───

def extract_text_color(image: np.ndarray, x_min: int, y_min: int, x_max: int, y_max: int) -> Tuple[int, int, int]:
    """
    从文字区域提取文字颜色。
    策略：K-Means 聚类找背景和前景，取与背景对比度最高的颜色。
    """
    region = image[y_min:y_max, x_min:x_max]
    if region.size == 0:
        return (0, 0, 0)

    region_rgb = cv2.cvtColor(region, cv2.COLOR_BGR2RGB)
    h, w = region_rgb.shape[:2]

    # 方法 1：K-Means 聚类
    result = _extract_color_kmeans(region_rgb)
    if result is not None:
        return result

    # 方法 2：背景/前景分离（Otsu）
    result = _extract_color_otsu(region_rgb)
    if result is not None:
        return result

    # 方法 3：回退到中心像素
    return tuple(int(c) for c in region_rgb[h // 2, w // 2])


def _extract_color_kmeans(region_rgb: np.ndarray, n_clusters: int = 3) -> Optional[Tuple[int, int, int]]:
    """K-Means 聚类提取文字颜色"""
    try:
        h, w = region_rgb.shape[:2]
        pixels = region_rgb.reshape(-1, 3).astype(np.float32)

        # 采样加速（超过 1000 像素时随机采样）
        if len(pixels) > 1000:
            indices = np.random.choice(len(pixels), 1000, replace=False)
            pixels = pixels[indices]

        # K-Means
        criteria = (cv2.TERM_CRITERIA_EPS + cv2.TERM_CRITERIA_MAX_ITER, 20, 1.0)
        _, labels, centers = cv2.kmeans(pixels, n_clusters, None, criteria, 3, cv2.KMEANS_PP_CENTERS)

        # 统计每个簇的像素数
        counts = np.bincount(labels.flatten(), minlength=n_clusters)

        # 背景 = 面积最大的簇
        bg_idx = np.argmax(counts)
        bg_color = centers[bg_idx]

        # 文字 = 与背景对比度最高的簇
        max_contrast = 0
        text_color = centers[0]
        for i in range(n_clusters):
            if i == bg_idx:
                continue
            contrast = np.sqrt(np.sum((centers[i] - bg_color) ** 2))
            if contrast > max_contrast:
                max_contrast = contrast
                text_color = centers[i]

        # 对比度太低说明没有明显文字
        if max_contrast < 30:
            return None

        return (int(text_color[0]), int(text_color[1]), int(text_color[2]))

    except Exception:
        return None


def _extract_color_otsu(region_rgb: np.ndarray) -> Optional[Tuple[int, int, int]]:
    """Otsu 二值化分离前景/背景"""
    try:
        gray = cv2.cvtColor(
            cv2.cvtColor(region_rgb, cv2.COLOR_RGB2BGR),
            cv2.COLOR_BGR2GRAY
        )
        _, binary = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)

        # 判断文字是亮色还是暗色（文字通常面积小于背景）
        fg_count = np.sum(binary == 0)
        bg_count = np.sum(binary == 255)

        if fg_count < bg_count:
            # 文字是暗色（黑字白底）
            text_mask = binary == 0
        else:
            # 文字是亮色（白字黑底）
            text_mask = binary == 255

        if text_mask.sum() == 0:
            return None

        text_pixels = region_rgb[text_mask]
        color = np.median(text_pixels, axis=0).astype(int)
        return (int(color[0]), int(color[1]), int(color[2]))

    except Exception:
        return None


# ─── 字号计算（二分查找） ───

def _get_font_path() -> Optional[str]:
    """获取系统字体路径"""
    font_paths = [
        "/System/Library/Fonts/PingFang.ttc",
        "/System/Library/Fonts/STHeiti Light.ttc",
        "/System/Library/Fonts/Helvetica.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    for fp in font_paths:
        if os.path.exists(fp):
            return fp
    return None


def fit_font_size(text: str, box_w: int, box_h: int, min_pt: int = 6, max_pt: int = 200) -> float:
    """
    基于文本框宽度估算字号，目标填充率约 75%。
    配合 SHAPE_TO_FIT_TEXT，文本框会自动调整高度适配文字。
    
    策略：二分查找，使文字渲染宽度 ≈ 框宽 × 0.75
    """
    font_path = _get_font_path()
    if not font_path:
        # 无字体时回退到高度估算
        return round(box_h * 0.75 * 0.75, 1)
    
    target_w = box_w * 0.75  # 目标填充宽度
    
    # 短文本（<=3字符）额外加高度约束，防止字号过大
    max_pt_limit = max_pt
    if len(text.strip()) <= 3:
        max_pt_limit = min(max_pt, box_h * 0.75 * 1.3)  # 允许略超框高
    
    low, high = min_pt, max_pt_limit
    best_pt = min_pt

    for _ in range(15):
        mid = (low + high) / 2
        try:
            font = ImageFont.truetype(font_path, int(mid))
            bbox = font.getbbox(text)
            text_w = bbox[2] - bbox[0]

            if text_w <= target_w:
                best_pt = mid
                low = mid + 0.5
            else:
                high = mid - 0.5
        except Exception:
            high = mid - 0.5

    return round(best_pt, 1)


# ─── 文字擦除 / 底图修复 ───

_lama_model = None


def get_lama_model(device="mps"):
    """懒加载 LaMa 模型"""
    global _lama_model
    if _lama_model is None:
        torch = _get_torch()
        from torch.hub import get_dir as get_torch_hub_dir

        if device == "mps" and torch.backends.mps.is_available():
            lama_device = torch.device("mps")
        else:
            lama_device = torch.device("cpu")

        model_path = Path(get_torch_hub_dir()) / "checkpoints" / "big-lama.pt"
        if not model_path.exists():
            from simple_lama_inpainting import SimpleLama
            SimpleLama(device=torch.device("cpu"))

        _lama_model = torch.jit.load(str(model_path), map_location=lama_device)
        _lama_model.eval()
        _lama_model = _lama_model.to(lama_device)

    return _lama_model


def _text_blocks_to_mask(shape: Tuple[int, int], blocks: List[TextBlock], dilate_px: int = 0) -> np.ndarray:
    h, w = shape
    mask = np.zeros((h, w), dtype=np.uint8)
    for block in blocks:
        x0 = max(0, block.x)
        y0 = max(0, block.y)
        x1 = min(w, block.x + block.w)
        y1 = min(h, block.y + block.h)
        if x1 > x0 and y1 > y0:
            mask[y0:y1, x0:x1] = 255

    if dilate_px > 0 and np.any(mask):
        kernel_size = dilate_px * 2 + 1
        kernel = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (kernel_size, kernel_size))
        mask = cv2.dilate(mask, kernel, iterations=1)
    return mask


def _opencv_inpaint(image_bgr: np.ndarray, mask: np.ndarray) -> np.ndarray:
    radius = int(os.environ.get("DECKLENS_INPAINT_RADIUS", "3"))
    return cv2.inpaint(image_bgr, mask, radius, cv2.INPAINT_TELEA)


def _block_border_pixels(region: np.ndarray) -> np.ndarray:
    h, w = region.shape[:2]
    border = max(2, min(h, w) // 8)
    top = region[:border, :, :].reshape(-1, 3)
    bottom = region[-border:, :, :].reshape(-1, 3)
    left = region[:, :border, :].reshape(-1, 3)
    right = region[:, -border:, :].reshape(-1, 3)
    return np.concatenate([top, bottom, left, right], axis=0)


def _build_text_pixel_mask(
    image_rgb: np.ndarray,
    blocks: List[TextBlock],
    min_diff: float = 24.0,
    mask_dilate: int = 1,
) -> np.ndarray:
    h, w = image_rgb.shape[:2]
    mask = np.zeros((h, w), dtype=np.uint8)

    for block in blocks:
        x0 = max(0, block.x)
        y0 = max(0, block.y)
        x1 = min(w, block.x + block.w)
        y1 = min(h, block.y + block.h)
        if x1 <= x0 or y1 <= y0:
            continue

        region = image_rgb[y0:y1, x0:x1]
        border = _block_border_pixels(region)
        bg = np.median(border.astype(np.float32), axis=0)
        diff = np.sqrt(np.sum((region.astype(np.float32) - bg) ** 2, axis=2)).astype(np.float32)
        diff_u8 = np.clip(diff, 0, 255).astype(np.uint8)
        otsu_threshold, _ = cv2.threshold(diff_u8, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
        threshold = max(float(otsu_threshold), min_diff)
        local = (diff >= threshold).astype(np.uint8)

        if int(local.sum()) < max(8, region.shape[0] * region.shape[1] * 0.003):
            inner_pad = max(1, min(region.shape[:2]) // 10)
            local = np.zeros(region.shape[:2], dtype=np.uint8)
            local[
                inner_pad : max(inner_pad + 1, region.shape[0] - inner_pad),
                inner_pad : max(inner_pad + 1, region.shape[1] - inner_pad),
            ] = 1

        mask[y0:y1, x0:x1] = np.maximum(mask[y0:y1, x0:x1], local * 255)

    if mask_dilate > 0 and np.any(mask):
        kernel_size = mask_dilate * 2 + 1
        kernel = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (kernel_size, kernel_size))
        mask = cv2.dilate(mask, kernel, iterations=1)
    return mask


def _component_bbox(component: np.ndarray, pad: int, width: int, height: int) -> Tuple[int, int, int, int]:
    ys, xs = np.where(component)
    x0 = max(0, int(xs.min()) - pad)
    y0 = max(0, int(ys.min()) - pad)
    x1 = min(width, int(xs.max()) + 1 + pad)
    y1 = min(height, int(ys.max()) + 1 + pad)
    return x0, y0, x1, y1


def _feather_alpha(component: np.ndarray, cover_dilate: int = 2, feather_blur: int = 5) -> np.ndarray:
    core = component.astype(np.uint8) * 255
    if cover_dilate > 0:
        kernel_size = cover_dilate * 2 + 1
        kernel = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (kernel_size, kernel_size))
        cover = cv2.dilate(core, kernel, iterations=1)
    else:
        cover = core

    if feather_blur <= 0:
        alpha = cover.astype(np.float32) / 255.0
    else:
        kernel_size = feather_blur * 2 + 1
        alpha = cv2.GaussianBlur(cover.astype(np.float32) / 255.0, (kernel_size, kernel_size), 0)
        if alpha.max() > 0:
            alpha = alpha / alpha.max()
        alpha = np.maximum(alpha, core.astype(np.float32) / 255.0)
    return np.clip(alpha, 0.0, 1.0)


def _local_mean_inpaint(image_bgr: np.ndarray, blocks: List[TextBlock]) -> np.ndarray:
    image_rgb = cv2.cvtColor(image_bgr, cv2.COLOR_BGR2RGB)
    h, w = image_rgb.shape[:2]
    text_mask = _build_text_pixel_mask(image_rgb, blocks)
    work = image_rgb.astype(np.float32).copy()
    binary = (text_mask > 0).astype(np.uint8)
    num_labels, labels, stats, _ = cv2.connectedComponentsWithStats(binary, connectivity=8)
    all_text = binary.astype(bool)
    sample_radius = int(os.environ.get("DECKLENS_LOCAL_MEAN_SAMPLE_RADIUS", "14"))
    min_component_area = int(os.environ.get("DECKLENS_LOCAL_MEAN_MIN_COMPONENT_AREA", "3"))

    for label_id in range(1, num_labels):
        area = int(stats[label_id, cv2.CC_STAT_AREA])
        if area < min_component_area:
            continue

        component = labels == label_id
        x0, y0, x1, y1 = _component_bbox(component, sample_radius + 9, w, h)
        local_component = component[y0:y1, x0:x1]
        local_all_text = all_text[y0:y1, x0:x1]
        local_image = image_rgb[y0:y1, x0:x1].astype(np.float32)

        kernel_size = sample_radius * 2 + 1
        kernel = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (kernel_size, kernel_size))
        ring_base = cv2.dilate(local_component.astype(np.uint8), kernel, iterations=1).astype(bool)
        sample_region = ring_base & ~local_all_text

        if int(sample_region.sum()) < 12:
            kernel_size = max(kernel_size + 8, 17)
            kernel = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (kernel_size, kernel_size))
            ring_base = cv2.dilate(local_component.astype(np.uint8), kernel, iterations=1).astype(bool)
            sample_region = ring_base & ~local_all_text

        if int(sample_region.sum()) == 0:
            fill_color = np.median(image_rgb.reshape(-1, 3).astype(np.float32), axis=0)
        else:
            fill_color = np.mean(local_image[sample_region], axis=0)

        alpha = _feather_alpha(local_component)[..., None]
        patch = np.ones_like(local_image) * fill_color.reshape(1, 1, 3)
        work[y0:y1, x0:x1] = patch * alpha + work[y0:y1, x0:x1] * (1.0 - alpha)

    result_rgb = np.clip(work, 0, 255).astype(np.uint8)
    return cv2.cvtColor(result_rgb, cv2.COLOR_RGB2BGR)


def _lama_inpaint(image_bgr: np.ndarray, mask: np.ndarray, device: str) -> np.ndarray:
    torch = _get_torch()
    model = get_lama_model(device)

    h, w = image_bgr.shape[:2]
    max_side = int(os.environ.get("DECKLENS_LAMA_MAX_SIDE", "1024"))
    scale = min(1.0, max(256, max_side) / max(h, w))
    work_bgr = image_bgr
    work_mask = mask
    if scale < 1.0:
        work_size = (max(1, int(round(w * scale))), max(1, int(round(h * scale))))
        work_bgr = cv2.resize(image_bgr, work_size, interpolation=cv2.INTER_AREA)
        work_mask = cv2.resize(mask, work_size, interpolation=cv2.INTER_NEAREST)

    rgb = cv2.cvtColor(work_bgr, cv2.COLOR_BGR2RGB).astype(np.float32) / 255.0
    mask_float = (work_mask > 0).astype(np.float32)
    image_tensor = torch.from_numpy(rgb).permute(2, 0, 1).unsqueeze(0)
    mask_tensor = torch.from_numpy(mask_float).unsqueeze(0).unsqueeze(0)

    try:
        model_device = next(model.parameters()).device
    except StopIteration:
        model_device = torch.device("mps" if device == "mps" and torch.backends.mps.is_available() else "cpu")
    image_tensor = image_tensor.to(model_device)
    mask_tensor = mask_tensor.to(model_device)

    with torch.no_grad():
        output = model(image_tensor, mask_tensor)

    output_rgb = output[0].detach().cpu().permute(1, 2, 0).numpy()
    output_rgb = np.clip(output_rgb * 255.0, 0, 255).astype(np.uint8)
    output_bgr = cv2.cvtColor(output_rgb, cv2.COLOR_RGB2BGR)
    if output_bgr.shape[:2] != (h, w):
        output_bgr = cv2.resize(output_bgr, (w, h), interpolation=cv2.INTER_CUBIC)
    return output_bgr


def remove_text_from_image(image_path: str, blocks: List[TextBlock], device="mps", backend: Optional[str] = None) -> Image.Image:
    """擦除文字，返回干净背景图"""
    image = cv2.imread(image_path)
    if image is None:
        raise ValueError(f"无法读取图片: {image_path}")
    h, w = image.shape[:2]

    backend = (backend or os.environ.get("DECKLENS_INPAINT_BACKEND", "lama")).strip().lower()
    hard_mask = _text_blocks_to_mask((h, w), blocks, dilate_px=0)

    if backend == "opencv":
        inpainted_bgr = _opencv_inpaint(image, hard_mask)
    elif backend == "local_mean":
        inpainted_bgr = _local_mean_inpaint(image, blocks)
    else:
        try:
            inpainted_bgr = _lama_inpaint(image, hard_mask, device=device)
        except Exception as exc:
            if not _env_truthy("DECKLENS_INPAINT_FALLBACK", True):
                raise
            print(f"  [INPAINT] LaMa 失败，回退本地均值: {exc}", flush=True)
            inpainted_bgr = _local_mean_inpaint(image, blocks)

    inpainted_rgb = cv2.cvtColor(inpainted_bgr, cv2.COLOR_BGR2RGB)
    return Image.fromarray(inpainted_rgb)


# ─── PPTX 生成 ───

PIXELS_TO_EMU = 9525  # 96 DPI: 1px = 914400/96 = 9525 EMU


def create_pptx(
    background_image: Image.Image,
    blocks: List[TextBlock],
    output_path: str,
):
    """合成 PPTX：背景图层 + 独立可编辑文字框"""
    img_w, img_h = background_image.size

    slide_width = img_w * PIXELS_TO_EMU
    slide_height = img_h * PIXELS_TO_EMU

    prs = Presentation()
    prs.slide_width = Emu(slide_width)
    prs.slide_height = Emu(slide_height)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # 背景图（JPG 压缩，减小体积）
    bg_path = output_path.replace(".pptx", "_bg.jpg")
    background_image.convert("RGB").save(bg_path, "JPEG", quality=85)
    slide.shapes.add_picture(
        bg_path, Emu(0), Emu(0),
        width=Emu(slide_width), height=Emu(slide_height)
    )

    # 文字框（每个检测框独立）
    for block in blocks:
        left = block.x * PIXELS_TO_EMU
        top = block.y * PIXELS_TO_EMU
        width = block.w * PIXELS_TO_EMU
        height = block.h * PIXELS_TO_EMU

        txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
        tf = txBox.text_frame

        # SHAPE_TO_FIT_TEXT: 文本框自动收缩适配文字
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        tf.word_wrap = False
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)

        # 垂直居中对齐
        from pptx.enum.text import PP_ALIGN
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT
        # 设置 bodyPr anchor="ctr"（垂直居中）
        from pptx.oxml.ns import qn as _qn_anchor
        body_pr = tf._txBody.bodyPr
        body_pr.set('anchor', 'ctr')

        p = tf.paragraphs[0]
        p.text = block.text

        # 行距 + 段前段后为0（消除字体行高带来的垂直偏移）
        from pptx.oxml.ns import qn
        from lxml import etree as _etree
        pPr = p._p.get_or_add_pPr()
        # 行距：纯数字/英文用 85%（微软雅黑上下留白大），中文用 100%
        import re as _re
        if _re.match(r'^[0-9A-Za-z\s\-_\.]+$', block.text.strip()):
            line_spacing_val = '85000'
        else:
            line_spacing_val = '100000'
        lnSpc = _etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = _etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', line_spacing_val)
        # 段前 0
        spcBef = _etree.SubElement(pPr, qn('a:spcBef'))
        spcPts = _etree.SubElement(spcBef, qn('a:spcPts'))
        spcPts.set('val', '0')
        # 段后 0
        spcAft = _etree.SubElement(pPr, qn('a:spcAft'))
        spcPts2 = _etree.SubElement(spcAft, qn('a:spcPts'))
        spcPts2.set('val', '0')

        # 字体属性
        run = p.runs[0]
        run.font.size = Pt(int(block.font_size_pt))
        r, g, b = block.color
        run.font.color.rgb = RGBColor(r, g, b)
        run.font.bold = block.bold

        # 字体设置：使用字体匹配器识别的字体
        run.font.name = block.font_name
        # 设置东亚字体
        try:
            from pptx.oxml.ns import qn
            from lxml import etree as _etree
            rPr = run._r.get_or_add_rPr()
            ea = _etree.SubElement(rPr, qn('a:ea'))
            ea.set('typeface', block.font_name)
        except Exception:
            pass

        # 文本框透明无边框
        txBox.fill.background()
        txBox.line.fill.background()

    prs.save(output_path)
    return output_path


# ─── PDF 拆页 ───

def pdf_to_images(pdf_path: str, dpi: int = 300, output_dir: Optional[str] = None) -> List[str]:
    """
    将 PDF 每页渲染为 PNG 图片，返回图片路径列表。
    """
    import fitz  # pymupdf

    pdf_path = str(Path(pdf_path).resolve())
    if output_dir is None:
        output_dir = str(Path(pdf_path).parent)

    doc = fitz.open(pdf_path)
    image_paths = []

    dpi = int(os.environ.get("DECKLENS_PDF_DPI", str(dpi)))
    base_zoom = dpi / 72.0

    for page_idx in range(len(doc)):
        page = doc[page_idx]
        rect = page.rect
        raw_w = int(rect.width * base_zoom)
        raw_h = int(rect.height * base_zoom)
        scale = _scale_for_limits(raw_w, raw_h)
        mat = fitz.Matrix(base_zoom * scale, base_zoom * scale)
        pix = page.get_pixmap(matrix=mat)
        img_path = str(Path(output_dir) / f"{Path(pdf_path).stem}_page{page_idx + 1:03d}.png")
        pix.save(img_path)
        image_paths.append(img_path)

    doc.close()
    return image_paths


# ─── 多图合并 PPTX（核心批量接口） ───

def _process_single_slide(
    image_path: str,
    prs: Presentation,
    device: str = "mps",
    expand_px: int = 3,
    bg_dir: str = "/tmp",
    slide_index: int = 0,
    progress_callback=None,
    decompose: bool = False,
    decompose_mode: str = "none",
    qwen_num_layers: int = 4,
    qwen_api_key: str = "",
    inpaint_backend: str = "lama",
):
    """
    处理单张图片，添加为 prs 中的一张 slide。
    """
    image_path = constrain_image_for_processing(str(Path(image_path).resolve()), output_dir=bg_dir)
    img = Image.open(image_path)
    img_w, img_h = img.size

    slide_width = img_w * PIXELS_TO_EMU
    slide_height = img_h * PIXELS_TO_EMU

    # 如果是第一张 slide，设置幻灯片尺寸；后续 slide 尺寸可能不同但 PPTX 只有一个全局尺寸
    # 取最大尺寸作为全局尺寸（后续 slide 居中放置）
    if slide_index == 0:
        prs.slide_width = Emu(slide_width)
        prs.slide_height = Emu(slide_height)
    else:
        # 更新为最大尺寸
        if slide_width > prs.slide_width:
            prs.slide_width = Emu(slide_width)
        if slide_height > prs.slide_height:
            prs.slide_height = Emu(slide_height)

    def report(msg):
        if progress_callback:
            progress_callback(msg)
        print(f"    {msg}")

    # OCR 检测
    report("OCR 检测文字...")
    t0 = time.time()
    blocks = detect_text_paddle(image_path, expand_px=expand_px)
    report(f"检测到 {len(blocks)} 个文字区域 ({time.time()-t0:.1f}s)")

    # 字体匹配
    if blocks:
        report("字体匹配...")
        t0 = time.time()
        blocks = match_fonts_for_blocks(image_path, blocks)
        report(f"字体匹配完成 ({time.time()-t0:.1f}s)")

    # 清理底图
    if blocks:
        backend_label = "本地均值" if inpaint_backend == "local_mean" else "LaMa"
        report(f"{backend_label} 清理底图...")
        t0 = time.time()
        background = remove_text_from_image(image_path, blocks, device=device, backend=inpaint_backend)
        report(f"背景生成完成 ({time.time()-t0:.1f}s)")
    else:
        background = img.convert("RGB")

    # 元素分层（可选）
    # decompose_mode: "none" / "sam" / "qwen"
    # 向后兼容: decompose=True 且 decompose_mode="none" 时视为 "sam"
    effective_mode = decompose_mode if decompose_mode in {"none", "sam", "qwen"} else "none"
    if effective_mode == "none" and decompose:
        effective_mode = "sam"

    try:
        qwen_num_layers = int(qwen_num_layers)
    except (TypeError, ValueError):
        qwen_num_layers = 4
    qwen_num_layers = max(3, min(8, qwen_num_layers))

    layers = []
    if effective_mode == "qwen":
        # Qwen-Image-Layered 云端分层
        report("AI 智能分层中（云端处理，预计 30-90s）...")
        t0 = time.time()
        # 先保存干净背景到临时文件供 API 上传
        qwen_input_path = os.path.join(bg_dir, f"slide_{slide_index:03d}_qwen_input.png")
        background.save(qwen_input_path, "PNG")
        layers = decompose_qwen(
            qwen_input_path,
            num_layers=qwen_num_layers,
            api_key=qwen_api_key,
        )
        if layers:
            report(f"AI 分层完成: {len(layers)} 层 ({time.time()-t0:.1f}s)")
        else:
            report(f"AI 分层失败，回退到 FastSAM 本地分层...")
            effective_mode = "sam"  # 回退

    if effective_mode == "sam" and not layers:
        report("FastSAM 元素分割中...")
        t0 = time.time()
        layers = decompose_background_sam(background, device=device)
        if layers:
            report(f"分割完成: {len(layers)-1} 个独立元素 ({time.time()-t0:.1f}s)")
        else:
            report(f"分割未产出有效元素，使用整张背景 ({time.time()-t0:.1f}s)")

    # 添加 slide
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    if layers:
        if effective_mode == "qwen":
            # Qwen 模式：每层是完整语义图层，直接按 bbox 放入
            layer_info = create_pptx_with_qwen_layers(
                layers, bg_dir=bg_dir, slide_index=slide_index,
                slide_w_px=img_w, slide_h_px=img_h,
            )
        else:
            # FastSAM 模式：底层整张 + 前景元素按连通域拆分
            layer_info = create_pptx_with_layers(
                layers, blocks, "", bg_dir=bg_dir, slide_index=slide_index
            )
        for path, x, y, w, h in layer_info:
            slide.shapes.add_picture(
                path,
                Emu(x * PIXELS_TO_EMU),
                Emu(y * PIXELS_TO_EMU),
                width=Emu(w * PIXELS_TO_EMU),
                height=Emu(h * PIXELS_TO_EMU),
            )
    else:
        # 单张背景图模式（原有逻辑）
        bg_path = str(Path(bg_dir) / f"slide_{slide_index:03d}_bg.jpg")
        background.convert("RGB").save(bg_path, "JPEG", quality=85)
        slide.shapes.add_picture(
            bg_path, Emu(0), Emu(0),
            width=Emu(slide_width), height=Emu(slide_height)
        )

    # 文字框
    for block in blocks:
        left = block.x * PIXELS_TO_EMU
        top = block.y * PIXELS_TO_EMU
        width = block.w * PIXELS_TO_EMU
        height = block.h * PIXELS_TO_EMU

        txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
        tf = txBox.text_frame

        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        tf.word_wrap = False
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)

        from pptx.enum.text import PP_ALIGN
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT
        from pptx.oxml.ns import qn as _qn_anchor
        body_pr = tf._txBody.bodyPr
        body_pr.set('anchor', 'ctr')

        p = tf.paragraphs[0]
        p.text = block.text

        from pptx.oxml.ns import qn
        from lxml import etree as _etree
        pPr = p._p.get_or_add_pPr()
        import re as _re
        if _re.match(r'^[0-9A-Za-z\s\-_\.]+$', block.text.strip()):
            line_spacing_val = '85000'
        else:
            line_spacing_val = '100000'
        lnSpc = _etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = _etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', line_spacing_val)
        spcBef = _etree.SubElement(pPr, qn('a:spcBef'))
        spcPts = _etree.SubElement(spcBef, qn('a:spcPts'))
        spcPts.set('val', '0')
        spcAft = _etree.SubElement(pPr, qn('a:spcAft'))
        spcPts2 = _etree.SubElement(spcAft, qn('a:spcPts'))
        spcPts2.set('val', '0')

        run = p.runs[0]
        run.font.size = Pt(int(block.font_size_pt))
        r, g, b = block.color
        run.font.color.rgb = RGBColor(r, g, b)
        run.font.bold = block.bold
        run.font.name = block.font_name
        try:
            rPr = run._r.get_or_add_rPr()
            ea = _etree.SubElement(rPr, qn('a:ea'))
            ea.set('typeface', block.font_name)
        except Exception:
            pass

        txBox.fill.background()
        txBox.line.fill.background()

    return slide


def images_to_pptx(
    image_paths: List[str],
    output_path: str,
    device: str = "mps",
    expand_px: int = 3,
    progress_callback=None,
    decompose: bool = False,
    decompose_mode: str = "none",
    qwen_num_layers: int = 4,
    qwen_api_key: str = "",
    inpaint_backend: str = "lama",
) -> str:
    """
    多张图片 → 一套 PPTX（每张图一个 slide）。
    progress_callback(current_index, total, filename, message)

    decompose_mode: "none" / "sam" / "qwen"
    """
    import tempfile
    bg_dir = tempfile.mkdtemp(prefix="decklens_bg_")

    prs = Presentation()
    total = len(image_paths)

    for idx, img_path in enumerate(image_paths):
        filename = Path(img_path).name

        def slide_progress(msg):
            if progress_callback:
                progress_callback(idx, total, filename, msg)

        slide_progress(f"处理中 ({idx+1}/{total})...")
        _process_single_slide(
            image_path=img_path,
            prs=prs,
            device=device,
            expand_px=expand_px,
            bg_dir=bg_dir,
            slide_index=idx,
            progress_callback=slide_progress,
            decompose=decompose,
            decompose_mode=decompose_mode,
            qwen_num_layers=qwen_num_layers,
            qwen_api_key=qwen_api_key,
            inpaint_backend=inpaint_backend,
        )
        slide_progress(f"完成")

    prs.save(output_path)

    # 清理临时背景图
    import shutil
    shutil.rmtree(bg_dir, ignore_errors=True)

    return output_path


# ─── 主流程（单图，保持向后兼容） ───

def image_to_pptx(
    image_path: str,
    output_path: Optional[str] = None,
    device: str = "mps",
    expand_px: int = 3,
    progress_callback=None,
    decompose: bool = False,
    decompose_mode: str = "none",
    qwen_num_layers: int = 4,
    qwen_api_key: str = "",
    inpaint_backend: str = "lama",
) -> str:
    """
    完整流程：单张图片 → 可编辑 PPTX（1 slide）
    """
    image_path = str(Path(image_path).resolve())

    if output_path is None:
        p = Path(image_path)
        output_path = str(p.parent / f"{p.stem}.pptx")

    def adapter(idx, total, filename, msg):
        if progress_callback:
            progress_callback(1, 1, msg)

    images_to_pptx(
        [image_path], output_path, device=device, expand_px=expand_px,
        progress_callback=adapter, decompose=decompose,
        decompose_mode=decompose_mode, qwen_num_layers=qwen_num_layers,
        qwen_api_key=qwen_api_key,
        inpaint_backend=inpaint_backend,
    )
    return output_path


# ─── 背景元素分割（FastSAM / SAM / OpenCV） ───

_sam_model = None
_sam_mask_generator = None
_fastsam_model = None


def release_cached_models():
    """Release heavyweight model singletons after a task to keep desktop RSS bounded."""
    global _paddle_ocr, _lama_model, _sam_model, _sam_mask_generator, _fastsam_model

    keep_ocr_default = sys.platform == "win32"
    if not _env_truthy("DECKLENS_KEEP_OCR_MODEL", default=keep_ocr_default):
        _paddle_ocr = None
    if _env_truthy("DECKLENS_RELEASE_OPTIONAL_MODELS", default=True):
        _lama_model = None
        _sam_model = None
        _sam_mask_generator = None
        _fastsam_model = None

    gc.collect()
    torch = sys.modules.get("torch")
    if torch is not None:
        if torch.cuda.is_available():
            torch.cuda.empty_cache()
        if hasattr(torch, "mps") and hasattr(torch.mps, "empty_cache"):
            torch.mps.empty_cache()

SAM_MODELS = {
    "vit_b": {
        "label": "ViT-B",
        "registry": "vit_b",
        "url": "https://dl.fbaipublicfiles.com/segment_anything/sam_vit_b_01ec64.pth",
        "path": os.path.expanduser("~/.cache/sam/sam_vit_b_01ec64.pth"),
        "min_bytes": 350_000_000,
    },
    "vit_h": {
        "label": "ViT-H",
        "registry": "vit_h",
        "url": "https://dl.fbaipublicfiles.com/segment_anything/sam_vit_h_4b8939.pth",
        "path": os.path.expanduser("~/.cache/sam/sam_vit_h_4b8939.pth"),
        "min_bytes": 2_000_000_000,
    },
}

# Backward-compatible names for callers/tests that patch these directly.
SAM_CHECKPOINT_URL = SAM_MODELS["vit_b"]["url"]
SAM_CHECKPOINT_PATH = SAM_MODELS["vit_b"]["path"]
SAM_CHECKPOINT_MIN_BYTES = SAM_MODELS["vit_b"]["min_bytes"]


def _sam_model_config() -> dict:
    model_name = os.environ.get("DECKLENS_SAM_MODEL", "vit_b").strip().lower()
    aliases = {"b": "vit_b", "base": "vit_b", "h": "vit_h", "huge": "vit_h"}
    model_name = aliases.get(model_name, model_name)
    if model_name not in SAM_MODELS:
        print(f"  [SAM] 未知模型 {model_name!r}，改用 vit_b", flush=True)
        model_name = "vit_b"

    config = dict(SAM_MODELS[model_name])
    if (
        SAM_CHECKPOINT_PATH != SAM_MODELS["vit_b"]["path"]
        or SAM_CHECKPOINT_URL != SAM_MODELS["vit_b"]["url"]
        or SAM_CHECKPOINT_MIN_BYTES != SAM_MODELS["vit_b"]["min_bytes"]
    ):
        config.update({
            "url": SAM_CHECKPOINT_URL,
            "path": SAM_CHECKPOINT_PATH,
            "min_bytes": SAM_CHECKPOINT_MIN_BYTES,
        })
    return config


def _sam_checkpoint_min_bytes(config: Optional[dict] = None) -> int:
    """Return the minimum plausible size for the configured SAM checkpoint."""
    config = config or _sam_model_config()
    try:
        return int(os.environ.get("DECKLENS_SAM_MIN_BYTES", config["min_bytes"]))
    except ValueError:
        return int(config["min_bytes"])


def _is_sam_checkpoint_complete(path: str, config: Optional[dict] = None) -> bool:
    """Detect partial/corrupt SAM checkpoint downloads before torch tries to load them."""
    return os.path.exists(path) and os.path.getsize(path) >= _sam_checkpoint_min_bytes(config)


def _expected_size_from_content_range(header: str) -> int:
    try:
        return int(header.rsplit("/", 1)[1])
    except (IndexError, ValueError, AttributeError):
        return 0


def _download_sam_checkpoint():
    """下载 SAM ViT-H checkpoint（如果不存在或已有文件不完整）"""
    config = _sam_model_config()
    checkpoint_path = config["path"]
    checkpoint_url = config["url"]

    if _is_sam_checkpoint_complete(checkpoint_path, config):
        return checkpoint_path

    tmp_path = f"{checkpoint_path}.part"
    if os.path.exists(checkpoint_path):
        size_mb = os.path.getsize(checkpoint_path) / 1024 / 1024
        print(f"  [SAM] 检测到不完整模型文件 ({size_mb:.0f}MB)，转为断点续传...", flush=True)
        os.replace(checkpoint_path, tmp_path)

    os.makedirs(os.path.dirname(checkpoint_path), exist_ok=True)

    import requests

    attempts = int(os.environ.get("DECKLENS_SAM_DOWNLOAD_RETRIES", "2"))
    read_timeout = float(os.environ.get("DECKLENS_SAM_DOWNLOAD_TIMEOUT", "20"))
    max_seconds = float(os.environ.get("DECKLENS_SAM_DOWNLOAD_MAX_SECONDS", "60"))
    started_at = time.time()
    expected_size = 0

    for attempt in range(1, attempts + 1):
        if time.time() - started_at > max_seconds:
            break

        existing_size = os.path.getsize(tmp_path) if os.path.exists(tmp_path) else 0
        headers = {"Range": f"bytes={existing_size}-"} if existing_size else {}
        resume_msg = f"，从 {existing_size / 1024 / 1024:.0f}MB 继续" if existing_size else ""
        print(
            f"  [SAM] 下载 {config['label']} 模型 ({attempt}/{attempts}{resume_msg})...",
            flush=True,
        )

        try:
            with requests.get(checkpoint_url, headers=headers, stream=True, timeout=(10, read_timeout)) as resp:
                resp.raise_for_status()

                if resp.status_code == 206:
                    expected_size = _expected_size_from_content_range(resp.headers.get("Content-Range", ""))
                    mode = "ab"
                else:
                    expected_size = int(resp.headers.get("Content-Length") or 0)
                    mode = "wb"

                with open(tmp_path, mode) as f:
                    for chunk in resp.iter_content(chunk_size=1024 * 1024):
                        if time.time() - started_at > max_seconds:
                            raise TimeoutError(f"SAM 模型下载超过 {max_seconds:.0f}s")
                        if chunk:
                            f.write(chunk)

            actual_size = os.path.getsize(tmp_path) if os.path.exists(tmp_path) else 0
            if expected_size and actual_size != expected_size:
                print(
                    f"  [SAM] 下载未完成: {actual_size / 1024 / 1024:.0f}MB / {expected_size / 1024 / 1024:.0f}MB",
                    flush=True,
                )
                continue

            if not _is_sam_checkpoint_complete(tmp_path, config):
                print(f"  [SAM] 下载文件过小: {actual_size / 1024 / 1024:.0f}MB", flush=True)
                continue

            os.replace(tmp_path, checkpoint_path)
            print(f"  [SAM] 模型下载完成", flush=True)
            return checkpoint_path
        except Exception as e:
            print(f"  [SAM] 下载中断: {e}", flush=True)
            if attempt == attempts:
                break

    actual_size = os.path.getsize(tmp_path) if os.path.exists(tmp_path) else 0
    expected_msg = f" / {expected_size / 1024 / 1024:.0f}MB" if expected_size else ""
    raise RuntimeError(f"SAM 模型下载失败，已保留断点文件: {actual_size / 1024 / 1024:.0f}MB{expected_msg}")


def get_sam_mask_generator(device="mps"):
    """懒加载 SAM + SamAutomaticMaskGenerator"""
    global _sam_model, _sam_mask_generator
    if not _env_truthy("DECKLENS_ENABLE_SAM", default=False):
        raise RuntimeError("本地 SAM 默认关闭以避免高内存峰值；如需启用请设置 DECKLENS_ENABLE_SAM=1")

    if _sam_mask_generator is None:
        torch = _get_torch()
        from segment_anything import sam_model_registry, SamAutomaticMaskGenerator

        config = _sam_model_config()
        checkpoint = _download_sam_checkpoint()

        if device == "mps" and torch.backends.mps.is_available():
            sam_device = torch.device("mps")
        else:
            sam_device = torch.device("cpu")

        print(f"  [SAM] 加载 {config['label']} 模型 (device={sam_device})...", flush=True)
        _sam_model = sam_model_registry[config["registry"]](checkpoint=checkpoint)
        _sam_model.to(sam_device)

        # MPS 不支持 float64，monkey-patch _process_batch 确保 float32
        if sam_device.type == "mps":
            _patch_sam_for_mps()

        _sam_mask_generator = SamAutomaticMaskGenerator(
            model=_sam_model,
            points_per_side=int(os.environ.get("DECKLENS_SAM_POINTS_PER_SIDE", "16")),
            pred_iou_thresh=0.86,
            stability_score_thresh=0.92,
            min_mask_region_area=500,
        )
        print(f"  [SAM] 模型加载完成", flush=True)

    return _sam_mask_generator


def _patch_sam_for_mps():
    """Monkey-patch SAM 的 automatic_mask_generator 以兼容 MPS (不支持 float64)"""
    import segment_anything.automatic_mask_generator as amg
    from segment_anything.utils.transforms import ResizeLongestSide

    # Patch apply_coords 返回 float32
    original_apply_coords = ResizeLongestSide.apply_coords

    def _patched_apply_coords(self, coords, original_size):
        result = original_apply_coords(self, coords, original_size)
        return result.astype(np.float32)

    ResizeLongestSide.apply_coords = _patched_apply_coords


def _fastsam_default_weights() -> str:
    env_path = os.environ.get("DECKLENS_FASTSAM_WEIGHTS")
    if env_path:
        return env_path

    data_dir = Path(os.environ.get("DECKLENS_DATA_DIR", Path(__file__).parent)).resolve()
    cache_path = data_dir / "models" / "ultralytics" / "FastSAM-s.pt"
    if cache_path.exists():
        return str(cache_path)

    dev_path = Path(__file__).parent / "test-materials" / "models" / "ultralytics" / "FastSAM-s.pt"
    if dev_path.exists():
        return str(dev_path)

    return "FastSAM-s.pt"


def _normalize_segmentation_mask(mask: np.ndarray, shape: Tuple[int, int]) -> np.ndarray:
    h, w = shape
    arr = np.asarray(mask)
    if arr.shape[:2] != (h, w):
        arr = cv2.resize(arr.astype(np.float32), (w, h), interpolation=cv2.INTER_NEAREST)
    if arr.dtype == bool:
        return arr
    return arr.astype(np.float32) > 0.5


def _bbox_from_segmentation(seg: np.ndarray) -> List[int]:
    ys, xs = np.where(seg)
    if len(xs) == 0 or len(ys) == 0:
        return [0, 0, 0, 0]
    x0 = int(xs.min())
    y0 = int(ys.min())
    x1 = int(xs.max())
    y1 = int(ys.max())
    return [x0, y0, x1 - x0 + 1, y1 - y0 + 1]


def get_fastsam_model():
    """Lazy-load FastSAM-s for local element layering."""
    global _fastsam_model
    if _fastsam_model is None:
        from ultralytics import FastSAM

        weights = _fastsam_default_weights()
        print(f"  [FastSAM] 加载模型: {weights}", flush=True)
        _fastsam_model = FastSAM(weights)
        print("  [FastSAM] 模型加载完成", flush=True)
    return _fastsam_model


def generate_background_fastsam_masks(image_rgb: np.ndarray, device: str = "cpu") -> List[dict]:
    """
    Generate SAM-like mask dictionaries with FastSAM-s.

    The downstream post-processing expects dicts with segmentation, area, and bbox.
    """
    model = get_fastsam_model()
    imgsz = int(os.environ.get("DECKLENS_FASTSAM_IMGSZ", "1024"))
    conf = float(os.environ.get("DECKLENS_FASTSAM_CONF", "0.25"))
    torch = sys.modules.get("torch")
    if torch is None:
        try:
            torch = _get_torch()
        except Exception:
            torch = None
    predict_device = "mps" if device == "mps" and torch is not None and torch.backends.mps.is_available() else "cpu"

    result = model.predict(
        image_rgb,
        imgsz=imgsz,
        conf=conf,
        retina_masks=True,
        verbose=False,
        device=predict_device,
    )[0]
    if result.masks is None:
        return []

    scores = result.boxes.conf.cpu().numpy().tolist() if result.boxes is not None else []
    boxes = result.boxes.xyxy.cpu().numpy().tolist() if result.boxes is not None else []
    masks = result.masks.data.cpu().numpy()
    out = []
    for idx, mask in enumerate(masks):
        seg = _normalize_segmentation_mask(mask, image_rgb.shape[:2])
        if idx < len(boxes):
            x0, y0, x1, y1 = boxes[idx]
            bbox = [int(round(x0)), int(round(y0)), int(round(x1 - x0)), int(round(y1 - y0))]
        else:
            bbox = _bbox_from_segmentation(seg)
        out.append(
            {
                "segmentation": seg,
                "area": int(seg.sum()),
                "bbox": bbox,
                "label": "fastsam",
                "score": round(float(scores[idx]), 4) if idx < len(scores) else None,
            }
        )
    print(f"  [FastSAM] 原始 mask 数量: {len(out)}", flush=True)
    return out


def segment_background_cv_masks(
    image_rgb: np.ndarray,
    min_area_ratio: float = 0.001,
    max_area_ratio: float = 0.80,
    bg_color_threshold: float = 30.0,
) -> list:
    """
    Lightweight fallback segmentation when SAM is unavailable.

    Returns SAM-like mask dictionaries: segmentation, area, bbox.
    """
    img_h, img_w = image_rgb.shape[:2]
    total_area = img_w * img_h

    corner_size = max(20, min(img_w, img_h) // 10)
    corners = np.concatenate([
        image_rgb[:corner_size, :corner_size].reshape(-1, 3),
        image_rgb[:corner_size, -corner_size:].reshape(-1, 3),
        image_rgb[-corner_size:, :corner_size].reshape(-1, 3),
        image_rgb[-corner_size:, -corner_size:].reshape(-1, 3),
    ], axis=0)
    bg_color = np.median(corners, axis=0)

    diff = np.sqrt(np.sum((image_rgb.astype(np.float32) - bg_color.astype(np.float32)) ** 2, axis=2))
    foreground = (diff >= bg_color_threshold).astype(np.uint8)

    kernel = np.ones((5, 5), np.uint8)
    foreground = cv2.morphologyEx(foreground, cv2.MORPH_OPEN, kernel)
    foreground = cv2.morphologyEx(foreground, cv2.MORPH_CLOSE, kernel, iterations=2)

    num_labels, labels, stats, _centroids = cv2.connectedComponentsWithStats(foreground, connectivity=8)
    masks = []
    for label_id in range(1, num_labels):
        x, y, w, h, area = stats[label_id]
        area_ratio = area / total_area
        if area_ratio < min_area_ratio or area_ratio > max_area_ratio:
            continue
        seg = labels == label_id
        masks.append({
            "segmentation": seg,
            "area": int(area),
            "bbox": [int(x), int(y), int(w), int(h)],
        })

    masks.sort(key=lambda m: m["area"], reverse=True)
    print(f"  [CV] 兜底分割生成 {len(masks)} 个候选元素", flush=True)
    return masks


def _dilate_binary_mask(mask: np.ndarray, radius_px: int = 5) -> np.ndarray:
    mask_u8 = (mask > 0).astype(np.uint8) * 255
    radius_px = max(0, int(radius_px))
    if radius_px <= 0:
        return mask_u8
    kernel = cv2.getStructuringElement(
        cv2.MORPH_ELLIPSE,
        (radius_px * 2 + 1, radius_px * 2 + 1),
    )
    return cv2.dilate(mask_u8, kernel, iterations=1)


def _telea_repair_rgb(image_rgb: np.ndarray, mask: np.ndarray, radius_px: int = 5) -> np.ndarray:
    repair_mask = _dilate_binary_mask(mask, radius_px=radius_px)
    if not np.any(repair_mask):
        return image_rgb.copy()
    image_bgr = cv2.cvtColor(image_rgb, cv2.COLOR_RGB2BGR)
    repaired_bgr = cv2.inpaint(image_bgr, repair_mask, 5, cv2.INPAINT_TELEA)
    return cv2.cvtColor(repaired_bgr, cv2.COLOR_BGR2RGB)


def create_deduped_rgba_layers(
    image_rgb: np.ndarray,
    segmentations: List[np.ndarray],
    repair_radius_px: int = 5,
    child_overlap_threshold: float = 0.85,
) -> List[Image.Image]:
    """
    Build editable element layers without duplicate pixels under movable elements.

    Layer 0 is the background with all selected foreground masks repaired away.
    For foreground layers, child masks that are fully inside a larger parent mask
    are repaired out of that parent layer, so moving the child does not reveal an
    identical copy still baked into the parent.
    """
    if image_rgb.ndim != 3 or image_rgb.shape[2] != 3:
        raise ValueError("image_rgb must be an RGB array")

    img_h, img_w = image_rgb.shape[:2]
    normalized = []
    for source_index, seg in enumerate(segmentations):
        seg_bool = np.asarray(seg).astype(bool)
        if seg_bool.shape != (img_h, img_w) or not seg_bool.any():
            continue
        normalized.append({
            "segmentation": seg_bool,
            "area": int(seg_bool.sum()),
            "source_index": source_index,
        })

    if not normalized:
        return [Image.fromarray(image_rgb).convert("RGBA")]

    # PPT shape stacking follows insertion order. Put larger parent/container
    # layers lower, then smaller child/detail layers above them.
    normalized.sort(key=lambda item: (-item["area"], item["source_index"]))

    union_mask = np.zeros((img_h, img_w), dtype=bool)
    for item in normalized:
        seg = item["segmentation"]
        union_mask = np.logical_or(union_mask, seg)

    background_rgb = _telea_repair_rgb(
        image_rgb,
        union_mask,
        radius_px=repair_radius_px,
    )
    layers = [Image.fromarray(background_rgb).convert("RGBA")]

    areas = [item["area"] for item in normalized]
    segment_masks = [item["segmentation"] for item in normalized]
    for idx, seg in enumerate(segment_masks):
        removal_mask = np.zeros((img_h, img_w), dtype=bool)
        for other_idx, other in enumerate(segment_masks):
            if other_idx == idx or areas[other_idx] >= areas[idx]:
                continue
            overlap = int(np.logical_and(seg, other).sum())
            if areas[other_idx] > 0 and overlap / areas[other_idx] >= child_overlap_threshold:
                removal_mask = np.logical_or(removal_mask, other)

        layer_rgb = image_rgb
        if removal_mask.any():
            layer_rgb = _telea_repair_rgb(
                image_rgb,
                np.logical_and(removal_mask, seg),
                radius_px=repair_radius_px,
            )

        elem_rgba = np.zeros((img_h, img_w, 4), dtype=np.uint8)
        elem_rgba[:, :, :3] = layer_rgb
        elem_rgba[:, :, 3] = (seg.astype(np.uint8) * 255)
        layers.append(Image.fromarray(elem_rgba, "RGBA"))

    return layers


def decompose_background_sam(
    background_image: Image.Image,
    device: str = "mps",
    timeout_sec: int = 120,
    min_area_ratio: float = 0.001,
    max_area_ratio: float = 0.80,
    bg_color_threshold: float = 30.0,
    overlap_iou_threshold: float = 0.7,
) -> List[Image.Image]:
    """
    用 FastSAM-s 将干净背景图分割为独立元素，失败时回退 OpenCV 轻量分割。

    参数:
        background_image: LaMa 擦除文字后的 RGB 图像
        device: 推理设备 ("mps" / "cpu")
        timeout_sec: 超时秒数
        min_area_ratio: 最小面积占比（低于此值视为噪声）
        max_area_ratio: 最大面积占比（高于此值视为背景）
        bg_color_threshold: 背景色差异阈值（低于此值视为背景的一部分）
        overlap_iou_threshold: 重叠 IoU 阈值（高于此值合并）

    返回: [layer_0 (背景), element_1, element_2, ...] — RGBA PIL Image 列表
        layer_0 是整张背景（全不透明），后续是独立元素。
        如果分割失败或无有效元素，返回空列表。
    """
    import concurrent.futures

    def _run_fastsam(rgb_image, img_w, img_h):
        """执行 FastSAM 推理 + 后处理"""
        image_np = np.array(rgb_image)  # (H, W, 3) uint8 RGB
        total_area = img_w * img_h

        masks = generate_background_fastsam_masks(image_np, device=device)

        if not masks:
            return []

        # --- 后处理 ---

        # 1. 面积过滤
        filtered = []
        for m in masks:
            area_ratio = m["area"] / total_area
            if area_ratio < min_area_ratio:
                continue  # 噪声碎片
            if area_ratio > max_area_ratio:
                continue  # 整张背景
            filtered.append(m)

        print(f"  [FastSAM] 面积过滤后: {len(filtered)}", flush=True)

        # 2. 背景色过滤
        # 计算全图背景色（取四角 + 中心区域的中位数）
        corner_size = max(20, min(img_w, img_h) // 10)
        corners = np.concatenate([
            image_np[:corner_size, :corner_size].reshape(-1, 3),
            image_np[:corner_size, -corner_size:].reshape(-1, 3),
            image_np[-corner_size:, :corner_size].reshape(-1, 3),
            image_np[-corner_size:, -corner_size:].reshape(-1, 3),
        ], axis=0)
        bg_color = np.median(corners, axis=0)

        color_filtered = []
        for m in filtered:
            seg = m["segmentation"]  # bool (H, W)
            mask_pixels = image_np[seg]
            mask_mean_color = np.mean(mask_pixels, axis=0)
            color_diff = np.sqrt(np.sum((mask_mean_color - bg_color) ** 2))
            if color_diff < bg_color_threshold:
                continue  # 与背景色太接近，跳过
            color_filtered.append(m)

        print(f"  [FastSAM] 背景色过滤后: {len(color_filtered)}", flush=True)

        # 3. 重叠合并（IoU > threshold 时保留面积更大的）
        # 按面积降序排列
        color_filtered.sort(key=lambda x: x["area"], reverse=True)
        keep = []
        for m in color_filtered:
            should_keep = True
            seg_m = m["segmentation"]
            for kept in keep:
                seg_k = kept["segmentation"]
                intersection = np.logical_and(seg_m, seg_k).sum()
                union = np.logical_or(seg_m, seg_k).sum()
                if union > 0 and intersection / union > overlap_iou_threshold:
                    should_keep = False
                    break
            if should_keep:
                keep.append(m)

        print(f"  [FastSAM] 去重后: {len(keep)}", flush=True)

        if not keep:
            del masks, filtered, color_filtered, keep
            return []

        layers = create_deduped_rgba_layers(
            image_np,
            [m["segmentation"] for m in keep],
        )

        del masks, filtered, color_filtered, keep
        return layers

    def _fallback_layers(rgb_image, img_w, img_h):
        image_np = np.array(rgb_image)
        masks = segment_background_cv_masks(
            image_np,
            min_area_ratio=min_area_ratio,
            max_area_ratio=max_area_ratio,
            bg_color_threshold=bg_color_threshold,
        )
        if not masks:
            return []

        layers = create_deduped_rgba_layers(
            image_np,
            [m["segmentation"] for m in masks],
        )
        del masks
        return layers

    try:
        rgb_image = background_image.convert("RGB")
        img_w, img_h = rgb_image.size

        if os.environ.get("DECKLENS_SEGMENT_BACKEND", "fastsam").strip().lower() == "opencv":
            print("  [FastSAM] 当前环境指定使用 OpenCV 兜底分割", flush=True)
            return _fallback_layers(rgb_image, img_w, img_h)

        with concurrent.futures.ThreadPoolExecutor(max_workers=1) as executor:
            future = executor.submit(_run_fastsam, rgb_image, img_w, img_h)
            try:
                result = future.result(timeout=timeout_sec)
                if result:
                    return result
                print("  [FastSAM] 未产出有效元素，改用 OpenCV 兜底分割", flush=True)
                return _fallback_layers(rgb_image, img_w, img_h)
            except concurrent.futures.TimeoutError:
                print(f"  [FastSAM] 分割超时 ({timeout_sec}s)，改用 OpenCV 兜底分割", flush=True)
                return _fallback_layers(rgb_image, img_w, img_h)

    except Exception as e:
        print(f"  [FastSAM] 分割失败: {e}", flush=True)
        import traceback
        traceback.print_exc()
        try:
            rgb_image = background_image.convert("RGB")
            img_w, img_h = rgb_image.size
            print("  [FastSAM] 改用 OpenCV 兜底分割", flush=True)
            return _fallback_layers(rgb_image, img_w, img_h)
        except Exception:
            return []


# ─── 连通域元素拆分 ───

@dataclass
class LayerElement:
    """一个从前景层拆分出的独立元素"""
    image: Image.Image  # RGBA 裁剪图
    x: int
    y: int
    w: int
    h: int


def split_layer_elements(
    layer: Image.Image,
    min_area: int = 100,
    erode_px: int = 0,
    alpha_threshold: int = 128,
) -> List[LayerElement]:
    """
    对 RGBA 前景层做 alpha 通道连通域分析，将不相连的区域各自裁出。
    
    参数:
        layer: RGBA PIL Image（前景层）
        min_area: 面积小于此值的碎片被过滤
        erode_px: 腐蚀像素数，可分离轻微接触的元素
        alpha_threshold: alpha 二值化阈值，低于此值视为透明
    
    返回: [LayerElement(image, x, y, w, h), ...]
    """
    layer_rgba = layer.convert("RGBA")
    alpha = np.array(layer_rgba.split()[3])

    # 二值化 alpha
    alpha_threshold = max(0, min(255, int(alpha_threshold)))
    binary = (alpha > alpha_threshold).astype(np.uint8)

    # 可选腐蚀（分离接触元素）
    if erode_px > 0:
        kernel = np.ones((erode_px * 2 + 1, erode_px * 2 + 1), np.uint8)
        binary = cv2.erode(binary, kernel, iterations=1)

    # 连通域分析
    num_labels, labels = cv2.connectedComponents(binary)

    elements = []
    layer_array = np.array(layer_rgba)

    for label_id in range(1, num_labels):  # 跳过背景 (0)
        mask = (labels == label_id)
        area = np.sum(mask)

        if area < min_area:
            continue

        # 计算 bbox
        rows = np.any(mask, axis=1)
        cols = np.any(mask, axis=0)
        y_min, y_max = np.where(rows)[0][[0, -1]]
        x_min, x_max = np.where(cols)[0][[0, -1]]

        w = x_max - x_min + 1
        h = y_max - y_min + 1

        # 裁剪元素（保留原始 RGBA，mask 外区域透明）
        elem_array = np.zeros((h, w, 4), dtype=np.uint8)
        crop_mask = mask[y_min:y_max + 1, x_min:x_max + 1]
        crop_rgba = layer_array[y_min:y_max + 1, x_min:x_max + 1]
        elem_array[crop_mask] = crop_rgba[crop_mask]

        elem_image = Image.fromarray(elem_array, "RGBA")
        elements.append(LayerElement(
            image=elem_image,
            x=int(x_min),
            y=int(y_min),
            w=int(w),
            h=int(h),
        ))

    return elements


def create_pptx_with_layers(
    layers: List[Image.Image],
    blocks: List[TextBlock],
    output_path: str,
    bg_dir: str = "/tmp",
    slide_index: int = 0,
):
    """
    用多图层模式生成 PPTX slide 内容。
    layers: RGBA 图层列表（从底到顶）
    blocks: 文字块列表
    
    底层整张放入，前景层拆分为独立元素按 bbox 定位。
    返回: [(path, x, y, w, h), ...] 每个元素的文件路径和位置
    """
    img_w, img_h = layers[0].size
    layer_info = []  # [(path, x, y, w, h)]
    elem_count = 0

    for i, layer in enumerate(layers):
        if i == 0:
            # 底层：整张放入，转 RGB 用 JPEG 压缩
            path = os.path.join(bg_dir, f"slide_{slide_index:03d}_layer_{i:02d}.jpg")
            layer.convert("RGB").save(path, "JPEG", quality=85)
            layer_info.append((path, 0, 0, img_w, img_h))
        else:
            # 前景层：连通域拆分为独立元素
            elements = split_layer_elements(layer, min_area=100, erode_px=0)
            for j, elem in enumerate(elements):
                path = os.path.join(bg_dir, f"slide_{slide_index:03d}_layer_{i:02d}_elem_{j:03d}.png")
                elem.image.save(path, "PNG")
                layer_info.append((path, elem.x, elem.y, elem.w, elem.h))
                elem_count += 1

    print(f"  [FastSAM] {len(layers)} 层, {elem_count} 个独立元素", flush=True)
    return layer_info


# ─── Qwen-Image-Layered 云端分层 ───

QWEN_TIMEOUT = 180  # 3 分钟超时（含排队）
DEFAULT_QWEN_FAL_MODEL_ID = "fal-ai/qwen-image-layered"
_FAL_CLIENT_LOCK = threading.Lock()


def decompose_qwen(
    image_path: str,
    num_layers: int = 4,
    resolution: int = 1024,
    num_inference_steps: int = 28,
    guidance_scale: float = 5.0,
    seed: Optional[int] = None,
    timeout: int = QWEN_TIMEOUT,
    api_key: str = "",
) -> List[Image.Image]:
    """
    调用 fal.ai Qwen-Image-Layered API 分解图层。

    需要本次任务传入 fal API Key，或设置环境变量 FAL_KEY。
    备选: 当 fal 不可用时回退到 HuggingFace Space (gradio_client)。

    参数:
        image_path: 输入图片路径（LaMa 去文字后的干净背景）
        num_layers: 分解层数 (3-8)
        resolution: 推理分辨率 (不直接暴露给 fal，fal 自动处理)
        num_inference_steps: 扩散步数
        guidance_scale: CFG 引导强度
        seed: 随机种子（None 则随机）
        timeout: API 超时秒数

    返回: RGBA PIL Image 列表（从底层到顶层）
        如果失败返回空列表。
    """
    try:
        num_layers = int(num_layers)
    except (TypeError, ValueError):
        num_layers = 4
    num_layers = max(3, min(8, num_layers))

    # 优先尝试 fal.ai
    layers = _decompose_via_fal(
        image_path,
        num_layers,
        num_inference_steps,
        guidance_scale,
        seed,
        api_key=api_key,
        timeout=timeout,
    )
    if layers:
        return layers

    # 回退到 HuggingFace Space
    print("  [Qwen] fal.ai 不可用，尝试 HuggingFace Space...", flush=True)
    return _decompose_via_hf_space(image_path, num_layers, resolution, num_inference_steps, guidance_scale, seed)


def _decompose_via_fal(
    image_path: str,
    num_layers: int = 4,
    num_inference_steps: int = 28,
    guidance_scale: float = 5.0,
    seed: Optional[int] = None,
    api_key: str = "",
    timeout: int = QWEN_TIMEOUT,
) -> List[Image.Image]:
    """通过 fal.ai API 调用 Qwen-Image-Layered"""
    try:
        import fal_client
    except ImportError:
        print("  [Qwen/fal] fal-client 未安装，请运行: pip install fal-client", flush=True)
        return []

    api_key = (api_key or "").strip()
    existing_key = os.environ.get("FAL_KEY")
    if not api_key and not existing_key:
        print("  [Qwen/fal] 未设置 FAL_KEY 环境变量", flush=True)
        return []

    try:
        import tempfile
        import shutil
        import urllib.request

        # fal 上传不支持中文文件名，复制到临时文件
        suffix = Path(image_path).suffix
        tmp_file = tempfile.NamedTemporaryFile(suffix=suffix, delete=False, prefix="qwen_input_")
        tmp_path = tmp_file.name
        tmp_file.close()
        shutil.copy2(image_path, tmp_path)
        previous_key = None

        try:
            with _FAL_CLIENT_LOCK:
                previous_key = os.environ.get("FAL_KEY")
                if api_key:
                    os.environ["FAL_KEY"] = api_key

                print(f"  [Qwen/fal] 上传图片...", flush=True)
                upload_url = fal_client.upload_file(tmp_path)

                print(
                    f"  [Qwen/fal] 开始分层 (model={DEFAULT_QWEN_FAL_MODEL_ID}, layers={num_layers}, steps={num_inference_steps})...",
                    flush=True,
                )
                arguments = {
                    "image_url": upload_url,
                    "num_layers": num_layers,
                    "num_inference_steps": num_inference_steps,
                    "guidance_scale": guidance_scale,
                    "output_format": "png",
                }
                if seed is not None:
                    arguments["seed"] = seed

                result = fal_client.subscribe(DEFAULT_QWEN_FAL_MODEL_ID, arguments=arguments, client_timeout=timeout)
        finally:
            if api_key:
                if previous_key is None:
                    os.environ.pop("FAL_KEY", None)
                else:
                    os.environ["FAL_KEY"] = previous_key
            os.unlink(tmp_path)

        # 下载图层
        layers = []
        for i, img_info in enumerate(result.get("images", [])):
            img_url = img_info["url"]
            layer_tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False, prefix=f"qwen_layer_{i}_")
            layer_path = layer_tmp.name
            layer_tmp.close()
            urllib.request.urlretrieve(img_url, layer_path)
            img = Image.open(layer_path).convert("RGBA")
            layers.append(img)
            os.unlink(layer_path)

        print(f"  [Qwen/fal] 分层完成: {len(layers)} 层", flush=True)
        return layers

    except Exception as e:
        print(f"  [Qwen/fal] API 调用失败: {e}", flush=True)
        import traceback
        traceback.print_exc()
        return []


def _decompose_via_hf_space(
    image_path: str,
    num_layers: int = 4,
    resolution: int = 1024,
    num_inference_steps: int = 50,
    guidance_scale: float = 4.0,
    seed: Optional[int] = None,
) -> List[Image.Image]:
    """通过 HuggingFace Space (gradio_client) 调用 Qwen-Image-Layered"""
    try:
        from gradio_client import Client
    except ImportError:
        print("  [Qwen/HF] gradio_client 未安装，请运行: pip install gradio_client", flush=True)
        return []

    try:
        print(f"  [Qwen/HF] 连接 HuggingFace Space...", flush=True)
        client = Client("Qwen/Qwen-Image-Layered")

        print(f"  [Qwen/HF] 开始分层 (layers={num_layers}, res={resolution})...", flush=True)
        arguments = {
            "image": image_path,
            "num_layers": num_layers,
            "resolution": resolution,
            "num_inference_steps": num_inference_steps,
            "true_cfg_scale": guidance_scale,
        }
        if seed is not None:
            arguments["seed"] = seed
        else:
            arguments["seed"] = 777

        result = client.predict(**arguments, api_name="/decompose")

        # result 可能是文件路径列表或单个目录
        if isinstance(result, (list, tuple)):
            layer_paths = result
        else:
            from pathlib import Path as _P
            result_path = _P(result)
            if result_path.is_dir():
                layer_paths = sorted(result_path.glob("*.png"))
            else:
                layer_paths = [result]

        layers = []
        for p in layer_paths:
            img = Image.open(str(p)).convert("RGBA")
            layers.append(img)

        print(f"  [Qwen/HF] 分层完成: {len(layers)} 层", flush=True)
        return layers

    except Exception as e:
        print(f"  [Qwen/HF] API 调用失败: {e}", flush=True)
        import traceback
        traceback.print_exc()
        return []


def create_pptx_with_qwen_layers(
    layers: List[Image.Image],
    bg_dir: str = "/tmp",
    slide_index: int = 0,
    slide_w_px: int = 1920,
    slide_h_px: int = 1080,
    min_element_area: int = 100,
    alpha_threshold: int = 16,
    background_coverage_threshold: float = 0.85,
) -> List[Tuple[str, int, int, int, int]]:
    """
    将 Qwen 输出的 RGBA 图层转为 PPTX 元素列表。

    Qwen 输出的每层是全尺寸 RGBA：
    - 第 0 层如果 alpha 覆盖率很高，作为底层背景，裁剪到有效区域后放入 PPTX。
      如果第 0 层也是稀疏透明元素，则同样进入连通域拆分。
    - 第 1~N 层先做 alpha 连通域分析，将同一图层里不相连的区域拆成
      独立 PNG，再按 bbox 放入 PPTX，避免多个不相连元素被锁在一个大图层里。

    返回: [(path, x, y, w, h), ...] 每个图层的文件路径和像素坐标
    """
    layer_info = []
    element_count = 0

    def _map_box(x: int, y: int, w: int, h: int, layer: Image.Image) -> Tuple[int, int, int, int]:
        scale_x = slide_w_px / layer.width
        scale_y = slide_h_px / layer.height
        return (
            int(round(x * scale_x)),
            int(round(y * scale_y)),
            max(1, int(round(w * scale_x))),
            max(1, int(round(h * scale_y))),
        )

    for i, layer in enumerate(layers):
        layer = layer.convert("RGBA")
        alpha = np.array(layer.split()[3])
        rows = np.any(alpha > 0, axis=1)
        cols = np.any(alpha > 0, axis=0)

        if not rows.any() or not cols.any():
            continue  # 全透明层，跳过

        opaque_area = int(np.count_nonzero(alpha > alpha_threshold))
        alpha_coverage = opaque_area / float(layer.width * layer.height)
        is_background_layer = i == 0 and alpha_coverage >= background_coverage_threshold

        if is_background_layer:
            # 底层（背景）：转 JPEG 减小体积
            y_min, y_max = np.where(rows)[0][[0, -1]]
            x_min, x_max = np.where(cols)[0][[0, -1]]
            cropped = layer.crop((x_min, y_min, x_max + 1, y_max + 1))
            crop_w = x_max - x_min + 1
            crop_h = y_max - y_min + 1

            path = os.path.join(bg_dir, f"slide_{slide_index:03d}_qwen_layer_{i:02d}.jpg")
            cropped.convert("RGB").save(path, "JPEG", quality=90)
            mapped_x, mapped_y, mapped_w, mapped_h = _map_box(
                int(x_min), int(y_min), int(crop_w), int(crop_h), layer
            )
            layer_info.append((path, mapped_x, mapped_y, mapped_w, mapped_h))
        else:
            # 前景层：先按 alpha 连通域拆成独立元素，再逐个写入 PPTX。
            elements = split_layer_elements(
                layer,
                min_area=min_element_area,
                erode_px=0,
                alpha_threshold=alpha_threshold,
            )
            for j, elem in enumerate(elements):
                path = os.path.join(bg_dir, f"slide_{slide_index:03d}_qwen_layer_{i:02d}_elem_{j:03d}.png")
                elem.image.save(path, "PNG")
                mapped_x, mapped_y, mapped_w, mapped_h = _map_box(elem.x, elem.y, elem.w, elem.h, layer)
                layer_info.append((path, mapped_x, mapped_y, mapped_w, mapped_h))
                element_count += 1

    print(
        f"  [Qwen] {len(layers)} 层 → {element_count} 个独立前景元素，"
        f"共写入 {len(layer_info)} 个 picture shape",
        flush=True,
    )
    return layer_info


if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("用法: python3 engine.py <图片路径|PDF路径> [输出路径.pptx]")
        print("      支持多个图片路径，空格分隔")
        sys.exit(1)

    inputs = sys.argv[1:-1] if len(sys.argv) > 2 and sys.argv[-1].endswith('.pptx') else sys.argv[1:]
    output = sys.argv[-1] if sys.argv[-1].endswith('.pptx') else None

    # 判断是否 PDF
    all_images = []
    for inp in inputs:
        if inp.lower().endswith('.pdf'):
            print(f"  PDF 拆页: {inp}")
            pages = pdf_to_images(inp)
            print(f"  拆出 {len(pages)} 页")
            all_images.extend(pages)
        else:
            all_images.append(inp)

    if output is None:
        p = Path(all_images[0])
        output = str(p.parent / f"{p.stem}.pptx")

    images_to_pptx(all_images, output)
