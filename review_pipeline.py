"""Agent-facing element review pipeline for DeckLens.

This module mirrors the app's background element preview flow, but writes the
intermediate state to disk so an Agent can inspect, merge, delete, and then
apply those decisions through the CLI.
"""

from __future__ import annotations

import colorsys
import json
import os
import pickle
import shutil
import tempfile
import time
from pathlib import Path
from typing import Callable, Iterable

import cv2
import numpy as np
from PIL import Image, ImageDraw, ImageFont

from engine import (
    PIXELS_TO_EMU,
    constrain_image_for_processing,
    create_deduped_rgba_layers,
    create_pptx_with_layers,
    detect_text_paddle,
    generate_background_fastsam_masks,
    match_fonts_for_blocks,
    pdf_to_images,
    release_cached_models,
    remove_text_from_image,
    segment_background_cv_masks,
)

SUPPORTED_INPUTS = {".png", ".jpg", ".jpeg", ".pdf"}
ProgressCallback = Callable[[int, int, str, str], None]


def release_memory() -> None:
    release_cached_models()
    import gc
    import sys

    gc.collect()
    torch = sys.modules.get("torch")
    if torch is not None:
        if torch.cuda.is_available():
            torch.cuda.empty_cache()
        if hasattr(torch, "mps") and hasattr(torch.mps, "empty_cache"):
            torch.mps.empty_cache()


def validate_review_inputs(inputs: Iterable[str]) -> list[Path]:
    paths: list[Path] = []
    for raw in inputs:
        path = Path(raw).expanduser().resolve()
        if not path.exists():
            raise FileNotFoundError(f"Input not found: {path}")
        if path.suffix.lower() not in SUPPORTED_INPUTS:
            raise ValueError(f"Unsupported input type: {path.suffix} ({path})")
        paths.append(path)
    if not paths:
        raise ValueError("At least one input file is required.")
    return paths


def expand_review_inputs(inputs: Iterable[str | Path], work_dir: str | Path) -> list[str]:
    work_dir = Path(work_dir)
    work_dir.mkdir(parents=True, exist_ok=True)
    image_paths: list[str] = []
    for raw in inputs:
        path = Path(raw).expanduser().resolve()
        if path.suffix.lower() == ".pdf":
            image_paths.extend(pdf_to_images(str(path), dpi=300, output_dir=str(work_dir)))
        else:
            image_paths.append(str(path))
    return image_paths


def _bbox_from_seg(seg: np.ndarray) -> list[int]:
    ys, xs = np.where(seg > 0)
    if len(xs) == 0 or len(ys) == 0:
        return [0, 0, 0, 0]
    x0, x1 = int(xs.min()), int(xs.max())
    y0, y1 = int(ys.min()), int(ys.max())
    return [x0, y0, x1 - x0 + 1, y1 - y0 + 1]


def _polygon_points(seg: np.ndarray) -> list[list[list[int]]]:
    seg_uint8 = (seg.astype(np.uint8) * 255)
    contours, _ = cv2.findContours(seg_uint8, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    polygons: list[list[list[int]]] = []
    for contour in contours:
        if len(contour) < 3:
            continue
        epsilon = 0.002 * cv2.arcLength(contour, True)
        approx = cv2.approxPolyDP(contour, epsilon, True)
        points = approx.reshape(-1, 2).astype(int).tolist()
        if len(points) >= 3:
            polygons.append(points)
    return polygons


def _mask_meta(mask: dict, mask_id: int, total_area: int) -> dict:
    hue = mask_id / 60
    r, g, b = colorsys.hsv_to_rgb(hue, 0.7, 0.9)
    bbox = [int(v) for v in mask.get("bbox") or _bbox_from_seg(mask["segmentation"])]
    area = int(mask["area"])
    x, y, w, h = bbox
    return {
        "id": mask_id,
        "area": area,
        "area_pct": round(area / max(total_area, 1) * 100, 2),
        "bbox": bbox,
        "center": [round(x + w / 2, 1), round(y + h / 2, 1)],
        "aspect_ratio": round(w / max(h, 1), 3),
        "color": [int(r * 255), int(g * 255), int(b * 255)],
        "polygons": _polygon_points(mask["segmentation"]),
    }


def _postprocess_masks(masks: list[dict], image_np: np.ndarray) -> list[dict]:
    img_h, img_w = image_np.shape[:2]
    total_area = img_w * img_h
    min_area_ratio = 0.001
    max_area_ratio = 0.80
    bg_color_threshold = 30.0
    overlap_iou_threshold = 0.7

    filtered = [m for m in masks if min_area_ratio <= m["area"] / total_area <= max_area_ratio]

    corner_size = max(20, min(img_w, img_h) // 10)
    corners = np.concatenate([
        image_np[:corner_size, :corner_size].reshape(-1, 3),
        image_np[:corner_size, -corner_size:].reshape(-1, 3),
        image_np[-corner_size:, :corner_size].reshape(-1, 3),
        image_np[-corner_size:, -corner_size:].reshape(-1, 3),
    ], axis=0)
    bg_color = np.median(corners, axis=0)

    color_filtered = []
    for m in filtered:
        seg = m["segmentation"]
        mask_mean_color = np.mean(image_np[seg], axis=0)
        color_diff = np.sqrt(np.sum((mask_mean_color - bg_color) ** 2))
        if color_diff >= bg_color_threshold:
            color_filtered.append(m)

    color_filtered.sort(key=lambda item: item["area"], reverse=True)
    keep = []
    for m in color_filtered:
        seg_m = m["segmentation"]
        should_keep = True
        for kept in keep:
            seg_k = kept["segmentation"]
            intersection = np.logical_and(seg_m, seg_k).sum()
            union = np.logical_or(seg_m, seg_k).sum()
            if union > 0 and intersection / union > overlap_iou_threshold:
                should_keep = False
                break
        if should_keep:
            keep.append(m)

    containment_filtered = []
    for i, m in enumerate(keep):
        seg_m = m["segmentation"]
        area_m = m["area"]
        total_covered = np.zeros_like(seg_m, dtype=bool)
        has_children = False
        for j, other in enumerate(keep):
            if i == j or other["area"] >= area_m:
                continue
            seg_other = other["segmentation"]
            overlap = np.logical_and(seg_m, seg_other).sum()
            if other["area"] > 0 and overlap / other["area"] > 0.9:
                total_covered = np.logical_or(total_covered, seg_other)
                has_children = True
        if has_children:
            coverage = total_covered.sum() / area_m if area_m > 0 else 0
            if coverage > 0.8:
                continue
        containment_filtered.append(m)
    return containment_filtered


def _segment_image(image_np: np.ndarray, device: str) -> list[dict]:
    try:
        if os.environ.get("DECKLENS_SEGMENT_BACKEND", "fastsam").strip().lower() == "opencv":
            raise RuntimeError("OpenCV fallback requested")
        masks = generate_background_fastsam_masks(image_np, device=device)
    except Exception as exc:
        print(f"  [FastSAM] review segmentation skipped/failed, using OpenCV fallback: {exc}", flush=True)
        masks = segment_background_cv_masks(image_np)
    return _postprocess_masks(masks, image_np)


def _draw_preview(bg_rgb: Image.Image, masks: list[dict], output_path: Path) -> None:
    img_w, img_h = bg_rgb.size
    vis = bg_rgb.copy().convert("RGBA")
    overlay = Image.new("RGBA", (img_w, img_h), (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)
    font = ImageFont.load_default()

    for midx, mask in enumerate(masks):
        hue = midx / max(len(masks), 1)
        r, g, b = colorsys.hsv_to_rgb(hue, 0.7, 0.9)
        color = (int(r * 255), int(g * 255), int(b * 255), 76)
        seg = mask["segmentation"]
        mask_img = Image.fromarray((seg * 255).astype(np.uint8), mode="L")
        colored = Image.new("RGBA", (img_w, img_h), color)
        overlay = Image.composite(colored, overlay, mask_img)

    vis = Image.alpha_composite(vis, overlay)
    draw = ImageDraw.Draw(vis)
    for midx, mask in enumerate(masks):
        x, y, w, h = [int(v) for v in mask["bbox"]]
        hue = midx / max(len(masks), 1)
        r, g, b = colorsys.hsv_to_rgb(hue, 0.72, 0.78)
        outline = (int(r * 255), int(g * 255), int(b * 255), 230)
        draw.rectangle([x, y, x + w, y + h], outline=outline, width=max(2, min(img_w, img_h) // 500))
        label = f"#{midx}"
        label_box = draw.textbbox((0, 0), label, font=font)
        label_w = label_box[2] - label_box[0] + 8
        label_h = label_box[3] - label_box[1] + 6
        label_x = max(0, min(x, img_w - label_w))
        label_y = max(0, y - label_h)
        draw.rounded_rectangle([label_x, label_y, label_x + label_w, label_y + label_h], radius=4, fill=(0, 0, 0, 210))
        draw.text((label_x + 4, label_y + 3), label, fill=(255, 255, 255, 255), font=font)

    vis.convert("RGB").save(output_path, "JPEG", quality=92)


def _save_mask_assets(bg_rgb: Image.Image, masks: list[dict], slide_dir: Path, root_dir: Path, max_size: int = 360) -> list[dict]:
    slide_dir.mkdir(parents=True, exist_ok=True)
    image_np = np.array(bg_rgb.convert("RGB"))
    assets: list[dict] = []
    for midx, mask in enumerate(masks):
        seg = mask["segmentation"].astype(bool)
        x, y, w, h = [int(v) for v in mask["bbox"]]
        pad = max(6, int(max(w, h) * 0.08))
        x0 = max(0, x - pad)
        y0 = max(0, y - pad)
        x1 = min(bg_rgb.width, x + w + pad)
        y1 = min(bg_rgb.height, y + h + pad)

        crop_rgb = image_np[y0:y1, x0:x1]
        crop_mask = seg[y0:y1, x0:x1]
        rgba = np.dstack([crop_rgb, (crop_mask.astype(np.uint8) * 255)])
        mask_path = slide_dir / f"mask_{midx:03d}.png"
        Image.fromarray(rgba, mode="RGBA").save(mask_path)

        context = Image.fromarray(crop_rgb, mode="RGB")
        context_path = slide_dir / f"mask_{midx:03d}_context.jpg"
        context.thumbnail((max_size, max_size), Image.Resampling.LANCZOS)
        canvas = Image.new("RGB", (max_size, max_size), "#f7f7f5")
        canvas.paste(context, ((max_size - context.width) // 2, (max_size - context.height) // 2))
        draw = ImageDraw.Draw(canvas)
        draw.rectangle([0, 0, max_size - 1, max_size - 1], outline="#c8c8c2", width=1)
        draw.rectangle([0, 0, 52, 24], fill="#111111")
        draw.text((7, 6), f"#{midx}", fill="#ffffff", font=ImageFont.load_default())
        canvas.save(context_path, "JPEG", quality=90)

        assets.append({
            "id": midx,
            "mask_image": str(mask_path.relative_to(root_dir)),
            "context_image": str(context_path.relative_to(root_dir)),
        })
    return assets


def _make_contact_sheet(context_paths: list[Path], output_path: Path, columns: int = 4) -> None:
    if not context_paths:
        return
    thumb = 220
    gap = 12
    rows = int(np.ceil(len(context_paths) / columns))
    sheet = Image.new("RGB", (columns * thumb + (columns + 1) * gap, rows * thumb + (rows + 1) * gap), "#f7f7f5")
    for idx, path in enumerate(context_paths):
        img = Image.open(path).convert("RGB")
        img.thumbnail((thumb, thumb), Image.Resampling.LANCZOS)
        x = gap + (idx % columns) * (thumb + gap)
        y = gap + (idx // columns) * (thumb + gap)
        sheet.paste(img, (x + (thumb - img.width) // 2, y + (thumb - img.height) // 2))
    sheet.save(output_path, "JPEG", quality=92)


def _bbox_overlap(a: list[int], b: list[int]) -> tuple[float, float, float]:
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    ix0 = max(ax, bx)
    iy0 = max(ay, by)
    ix1 = min(ax + aw, bx + bw)
    iy1 = min(ay + ah, by + bh)
    iw = max(0, ix1 - ix0)
    ih = max(0, iy1 - iy0)
    inter = iw * ih
    area_a = max(1, aw * ah)
    area_b = max(1, bw * bh)
    union = area_a + area_b - inter
    return inter / max(1, union), inter / min(area_a, area_b), inter / max(area_a, area_b)


def _suggest_merge_groups(mask_meta: list[dict], img_w: int, img_h: int) -> list[dict]:
    if len(mask_meta) < 2:
        return []
    parent = list(range(len(mask_meta)))

    def find(i: int) -> int:
        while parent[i] != i:
            parent[i] = parent[parent[i]]
            i = parent[i]
        return i

    def union(a: int, b: int) -> None:
        ra, rb = find(a), find(b)
        if ra != rb:
            parent[rb] = ra

    reasons: dict[tuple[int, int], str] = {}
    diag = max(1.0, (img_w ** 2 + img_h ** 2) ** 0.5)
    for i, a in enumerate(mask_meta):
        ax, ay, aw, ah = a["bbox"]
        acx, acy = a["center"]
        for j in range(i + 1, len(mask_meta)):
            b = mask_meta[j]
            bx, by, bw, bh = b["bbox"]
            bcx, bcy = b["center"]
            iou, contained_small, contained_large = _bbox_overlap(a["bbox"], b["bbox"])
            center_dist = ((acx - bcx) ** 2 + (acy - bcy) ** 2) ** 0.5 / diag
            gap_x = max(0, max(ax, bx) - min(ax + aw, bx + bw))
            gap_y = max(0, max(ay, by) - min(ay + ah, by + bh))
            near = gap_x <= max(10, min(aw, bw) * 0.18) and gap_y <= max(10, min(ah, bh) * 0.18)
            similar_scale = max(aw * ah, bw * bh) / max(1, min(aw * ah, bw * bh)) < 12
            if contained_small > 0.72:
                union(i, j)
                reasons[(i, j)] = "overlapping_or_nested_fragments"
            elif iou > 0.18:
                union(i, j)
                reasons[(i, j)] = "overlapping_fragments"
            elif near and similar_scale and center_dist < 0.18:
                union(i, j)
                reasons[(i, j)] = "nearby_similar_fragments"
            elif contained_large > 0.24 and center_dist < 0.10:
                union(i, j)
                reasons[(i, j)] = "same_visual_region"

    grouped: dict[int, list[int]] = {}
    for idx in range(len(mask_meta)):
        grouped.setdefault(find(idx), []).append(idx)

    suggestions = []
    for group in grouped.values():
        if len(group) < 2:
            continue
        xs = [mask_meta[i]["bbox"][0] for i in group]
        ys = [mask_meta[i]["bbox"][1] for i in group]
        x2 = [mask_meta[i]["bbox"][0] + mask_meta[i]["bbox"][2] for i in group]
        y2 = [mask_meta[i]["bbox"][1] + mask_meta[i]["bbox"][3] for i in group]
        group_reasons = sorted({
            reason
            for (a, b), reason in reasons.items()
            if a in group and b in group
        }) or ["likely_same_semantic_object"]
        suggestions.append({
            "ids": group,
            "bbox": [min(xs), min(ys), max(x2) - min(xs), max(y2) - min(ys)],
            "reasons": group_reasons,
            "agent_instruction": "Review visually. Merge if these masks belong to one photo, screenshot, icon, card, chart, or decoration.",
        })
    suggestions.sort(key=lambda item: (len(item["ids"]), item["bbox"][2] * item["bbox"][3]), reverse=True)
    return suggestions


def create_review(
    inputs: Iterable[str],
    review_dir: str,
    *,
    device: str = "cpu",
    inpaint_backend: str = "lama",
    progress_callback: ProgressCallback | None = None,
) -> dict:
    input_paths = validate_review_inputs(inputs)
    review_root = Path(review_dir).expanduser().resolve()
    if review_root.exists() and any(review_root.iterdir()):
        raise FileExistsError(f"Review directory already exists and is not empty: {review_root}")
    review_root.mkdir(parents=True, exist_ok=True)
    pages_dir = review_root / "pages"
    pages_dir.mkdir(exist_ok=True)

    expanded_dir = review_root / "expanded_inputs"
    image_paths = expand_review_inputs(input_paths, expanded_dir)
    slides = []
    started_at = time.time()

    for idx, image_path in enumerate(image_paths):
        progress_callback and progress_callback(idx, len(image_paths), image_path, "reading text")
        slide_dir = pages_dir / f"slide_{idx:03d}"
        slide_dir.mkdir(parents=True, exist_ok=True)
        constrained_path = constrain_image_for_processing(image_path, output_dir=str(slide_dir))
        original_image = Image.open(constrained_path).convert("RGB")
        original_copy = slide_dir / "original.png"
        original_image.save(original_copy)

        blocks = detect_text_paddle(constrained_path)
        if blocks:
            blocks = match_fonts_for_blocks(constrained_path, blocks)

        progress_callback and progress_callback(idx, len(image_paths), image_path, "cleaning background")
        if blocks:
            background = remove_text_from_image(constrained_path, blocks, device=device, backend=inpaint_backend)
        else:
            background = original_image
        bg_clean_path = slide_dir / "background_clean.png"
        background.convert("RGB").save(bg_clean_path)

        progress_callback and progress_callback(idx, len(image_paths), image_path, "segmenting elements")
        bg_rgb = background.convert("RGB")
        bg_np = np.array(bg_rgb)
        img_w, img_h = bg_rgb.size
        total_area = img_w * img_h
        masks = _segment_image(bg_np, device=device)

        seg_arrays = {f"mask_{midx}": mask["segmentation"].astype(np.uint8) for midx, mask in enumerate(masks)}
        masks_npz = slide_dir / "masks.npz"
        np.savez_compressed(masks_npz, **seg_arrays)
        with open(slide_dir / "blocks.pkl", "wb") as file:
            pickle.dump(blocks, file)

        for midx, mask in enumerate(masks):
            mask["bbox"] = [int(v) for v in mask.get("bbox") or _bbox_from_seg(mask["segmentation"])]

        preview_path = slide_dir / "preview_numbered.jpg"
        _draw_preview(bg_rgb, masks, preview_path)
        mask_assets = _save_mask_assets(bg_rgb, masks, slide_dir / "masks", review_root)
        contact_sheet = slide_dir / "contact_sheet.jpg"
        _make_contact_sheet([review_root / asset["context_image"] for asset in mask_assets], contact_sheet)

        mask_meta = []
        for midx, mask in enumerate(masks):
            meta = _mask_meta(mask, midx, total_area)
            asset = next((item for item in mask_assets if item["id"] == midx), None)
            if asset:
                meta.update({
                    "mask_image": asset["mask_image"],
                    "context_image": asset["context_image"],
                })
            mask_meta.append(meta)

        slides.append({
            "slide_index": idx,
            "source_image": str(Path(image_path).resolve()),
            "original_image": str(original_copy.relative_to(review_root)),
            "bg_clean": str(bg_clean_path.relative_to(review_root)),
            "preview_image": str(preview_path.relative_to(review_root)),
            "contact_sheet": str(contact_sheet.relative_to(review_root)) if contact_sheet.exists() else None,
            "masks_file": str(masks_npz.relative_to(review_root)),
            "blocks_file": str((slide_dir / "blocks.pkl").relative_to(review_root)),
            "img_width": img_w,
            "img_height": img_h,
            "blocks_count": len(blocks),
            "masks": mask_meta,
            "suggested_merge_groups": _suggest_merge_groups(mask_meta, img_w, img_h),
        })
        release_memory()

    decision_template = {
        "slides": [
            {
                "slide_index": slide["slide_index"],
                "merge": [group["ids"] for group in slide["suggested_merge_groups"]],
                "delete": [],
                "keep": [mask["id"] for mask in slide["masks"]],
            }
            for slide in slides
        ],
        "notes": "Edit merge/delete/keep after visual review. Remove suggested groups that are not one semantic object.",
    }
    (review_root / "decision.template.json").write_text(json.dumps(decision_template, ensure_ascii=False, indent=2), encoding="utf-8")

    manifest = {
        "schema": "decklens.review.v1",
        "created_at": time.strftime("%Y-%m-%dT%H:%M:%S%z"),
        "review_dir": str(review_root),
        "inputs": [str(path) for path in input_paths],
        "inpaint_backend": inpaint_backend,
        "device": device,
        "total_slides": len(slides),
        "elapsed_sec": round(time.time() - started_at, 2),
        "slides": slides,
        "decision_template": "decision.template.json",
        "agent_checklist": [
            "Open each preview_image and contact_sheet before deciding.",
            "Merge fragments that are one photo/screenshot/icon/card/chart/decoration.",
            "Delete masks that will be replaced by SVG, icon, shape, or a merged bitmap later.",
            "Keep masks only when they are independent semantic objects.",
        ],
    }
    manifest_path = review_root / "manifest.json"
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    return manifest


def _decision_for_slide(decision: dict, slide_index: int) -> dict | None:
    for slide in decision.get("slides", []):
        if int(slide.get("slide_index", -1)) == slide_index:
            return slide
    return None


def _validate_decision(slide: dict, decision: dict | None) -> tuple[set[int], list[list[int]], set[int]]:
    valid_ids = {int(mask["id"]) for mask in slide.get("masks", [])}
    if decision is None:
        return set(valid_ids), [], set()

    delete_ids = {int(item) for item in decision.get("delete", [])}
    merge_groups = [[int(item) for item in group] for group in decision.get("merge", [])]
    if "keep" in decision:
        keep_ids = {int(item) for item in decision.get("keep", [])}
    else:
        merged = {item for group in merge_groups for item in group}
        keep_ids = valid_ids - delete_ids - merged

    referenced = set(keep_ids) | delete_ids | {item for group in merge_groups for item in group}
    invalid = sorted(referenced - valid_ids)
    if invalid:
        raise ValueError(f"Decision references invalid mask ids on slide {slide['slide_index']}: {invalid}")
    return keep_ids, merge_groups, delete_ids


def apply_review(manifest_path: str, decision_path: str, output_path: str) -> dict:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.util import Emu, Pt
    from lxml import etree as _etree

    manifest_file = Path(manifest_path).expanduser().resolve()
    review_root = manifest_file.parent
    manifest = json.loads(manifest_file.read_text(encoding="utf-8"))
    decision = json.loads(Path(decision_path).expanduser().resolve().read_text(encoding="utf-8"))
    output = Path(output_path).expanduser().resolve()
    output.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    summary = {"slides": [], "output": str(output)}

    for slide_info in manifest["slides"]:
        idx = int(slide_info["slide_index"])
        img_w = int(slide_info["img_width"])
        img_h = int(slide_info["img_height"])
        slide_width = img_w * PIXELS_TO_EMU
        slide_height = img_h * PIXELS_TO_EMU
        if idx == 0:
            prs.slide_width = Emu(slide_width)
            prs.slide_height = Emu(slide_height)
        else:
            if slide_width > prs.slide_width:
                prs.slide_width = Emu(slide_width)
            if slide_height > prs.slide_height:
                prs.slide_height = Emu(slide_height)

        slide_decision = _decision_for_slide(decision, idx)
        keep_ids, merge_groups, delete_ids = _validate_decision(slide_info, slide_decision)
        masks_npz_path = review_root / slide_info["masks_file"]
        blocks_path = review_root / slide_info["blocks_file"]
        bg_path = review_root / slide_info["bg_clean"]

        with np.load(masks_npz_path) as masks_data:
            with open(blocks_path, "rb") as file:
                blocks = pickle.load(file)
            bg_image = Image.open(bg_path).convert("RGB")

            final_segs = []
            merged_ids = set()
            applied_merge_groups = []
            for group in merge_groups:
                merged_seg = np.zeros((img_h, img_w), dtype=np.uint8)
                valid_group = []
                for mid in group:
                    key = f"mask_{mid}"
                    if key in masks_data:
                        merged_seg = np.logical_or(merged_seg, masks_data[key]).astype(np.uint8)
                        merged_ids.add(mid)
                        valid_group.append(mid)
                if merged_seg.any():
                    final_segs.append(merged_seg)
                    applied_merge_groups.append(valid_group)

            for mid in sorted(keep_ids):
                if mid in merged_ids or mid in delete_ids:
                    continue
                key = f"mask_{mid}"
                if key in masks_data:
                    final_segs.append(masks_data[key].copy())

        bg_np = np.array(bg_image)
        layers = create_deduped_rgba_layers(bg_np, final_segs)
        bg_dir = tempfile.mkdtemp(prefix="decklens_review_apply_")
        try:
            layer_info = create_pptx_with_layers(layers, blocks, "", bg_dir=bg_dir, slide_index=idx)
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)

            for path, x, y, w, h in layer_info:
                slide.shapes.add_picture(
                    path,
                    Emu(x * PIXELS_TO_EMU),
                    Emu(y * PIXELS_TO_EMU),
                    width=Emu(w * PIXELS_TO_EMU),
                    height=Emu(h * PIXELS_TO_EMU),
                )

            for block in blocks:
                left = block.x * PIXELS_TO_EMU
                top = block.y * PIXELS_TO_EMU
                width = block.w * PIXELS_TO_EMU
                height = block.h * PIXELS_TO_EMU
                tx_box = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
                tf = tx_box.text_frame
                tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                tf.word_wrap = False
                tf.margin_left = Emu(0)
                tf.margin_right = Emu(0)
                tf.margin_top = Emu(0)
                tf.margin_bottom = Emu(0)
                tf.paragraphs[0].alignment = PP_ALIGN.LEFT
                tf._txBody.bodyPr.set("anchor", "ctr")

                p = tf.paragraphs[0]
                p.text = block.text
                p_pr = p._p.get_or_add_pPr()
                line_spacing_val = "85000" if __import__("re").match(r"^[0-9A-Za-z\s\-_\.]+$", block.text.strip()) else "100000"
                ln_spc = _etree.SubElement(p_pr, qn("a:lnSpc"))
                spc_pct = _etree.SubElement(ln_spc, qn("a:spcPct"))
                spc_pct.set("val", line_spacing_val)
                for name in ("spcBef", "spcAft"):
                    spc = _etree.SubElement(p_pr, qn(f"a:{name}"))
                    spc_pts = _etree.SubElement(spc, qn("a:spcPts"))
                    spc_pts.set("val", "0")

                run = p.runs[0]
                run.font.size = Pt(int(block.font_size_pt))
                r, g, b = block.color
                run.font.color.rgb = RGBColor(r, g, b)
                run.font.bold = block.bold
                run.font.name = block.font_name
                try:
                    r_pr = run._r.get_or_add_rPr()
                    ea = _etree.SubElement(r_pr, qn("a:ea"))
                    ea.set("typeface", block.font_name)
                except Exception:
                    pass
                tx_box.fill.background()
                tx_box.line.fill.background()
        finally:
            shutil.rmtree(bg_dir, ignore_errors=True)

        summary["slides"].append({
            "slide_index": idx,
            "kept": sorted(keep_ids - merged_ids - delete_ids),
            "merged": applied_merge_groups,
            "deleted": sorted(delete_ids),
            "output_picture_layers": len(final_segs) + 1,
            "text_blocks": len(blocks),
        })
        release_memory()

    prs.save(output)
    return summary
