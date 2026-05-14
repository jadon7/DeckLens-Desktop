"""
app.py — DeckLens Web UI
Flask 应用，提供拖拽上传 + 转换 + 下载功能
支持多张图片批量转换为一套 PPTX；PDF 作为页面图片导入来源

v7: 异步处理 + 实时进度轮询
v8: FastSAM 元素分割 + 交互式合并预览
v9: Qwen-Image-Layered 云端分层
"""

import os
import uuid
import time
import threading
import tempfile
import json
import gc
import sys
import shutil
import re
from pathlib import Path

import numpy as np
import cv2
from PIL import Image
from flask import Flask, request, jsonify, send_file, render_template

from engine import image_to_pptx, images_to_pptx, pdf_to_images, detect_text_paddle, remove_text_from_image, generate_background_fastsam_masks, decompose_background_sam, segment_background_cv_masks, constrain_image_for_processing, match_fonts_for_blocks, release_cached_models

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = int(os.environ.get("DECKLENS_MAX_UPLOAD_MB", "100")) * 1024 * 1024

BASE_DIR = Path(__file__).parent
DATA_DIR = Path(os.environ.get("DECKLENS_DATA_DIR", BASE_DIR)).resolve()
UPLOAD_DIR = DATA_DIR / "uploads"
OUTPUT_DIR = DATA_DIR / "outputs"
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
DEFAULT_DEVICE = os.environ.get("DECKLENS_DEVICE", "cpu")
INPAINT_BACKENDS = {"lama", "local_mean"}
INPAINT_BACKEND_LABELS = {
    "lama": "LaMa",
    "local_mean": "本地均值",
}
OUTPUT_MODE_LABELS = {
    "none": "标准还原",
    "sam": "元素分层",
    "qwen": "AI智能分层",
}
# 任务状态存储（简单内存版）
tasks = {}


def release_memory():
    """Best-effort cleanup after heavyweight OCR/inpainting/segmentation steps."""
    release_cached_models()
    gc.collect()
    torch = sys.modules.get("torch")
    if torch is not None:
        if torch.cuda.is_available():
            torch.cuda.empty_cache()
        if hasattr(torch, "mps") and hasattr(torch.mps, "empty_cache"):
            torch.mps.empty_cache()


def compact_finished_task(task):
    """Drop preview-only data after output is ready; download only needs output path."""
    for key in ("preview_data", "all_images", "user_decision"):
        task.pop(key, None)


def original_upload_name(saved_path: str) -> str:
    name = Path(saved_path).name
    return name.split("_", 1)[1] if "_" in name else name


def safe_filename_part(value: str, fallback: str = "DeckLens") -> str:
    stem = Path(value or fallback).stem.strip() or fallback
    stem = re.sub(r'[\\/:*?"<>|\x00-\x1f]+', "-", stem)
    stem = re.sub(r"\s+", " ", stem).strip(" .-_")
    return (stem or fallback)[:64]


def allocate_output_path(saved_paths: list, decompose_mode: str) -> str:
    first_name = original_upload_name(saved_paths[0]) if saved_paths else "DeckLens"
    base = safe_filename_part(first_name)
    if len(saved_paths) > 1:
        base = f"{base} 等{len(saved_paths)}个文件"
    mode = OUTPUT_MODE_LABELS.get(decompose_mode, OUTPUT_MODE_LABELS["none"])
    timestamp = time.strftime("%Y%m%d-%H%M%S")
    filename = f"{base}_{mode}_{timestamp}.pptx"
    candidate = OUTPUT_DIR / filename
    suffix = 2
    while candidate.exists():
        candidate = OUTPUT_DIR / f"{base}_{mode}_{timestamp}-{suffix}.pptx"
        suffix += 1
    return str(candidate)


@app.route("/")
def index():
    return render_template("index.html")


@app.route("/healthz")
def healthz():
    return jsonify({"ok": True, "service": "DeckLens"})


@app.route("/api/convert", methods=["POST"])
def convert():
    """
    批量转换接口（异步）。
    接收一个或多个文件（field name: files），立即返回 task_id。
    前端通过 /api/status/<task_id> 轮询进度。
    """
    uploaded_files = request.files.getlist("files")
    if not uploaded_files:
        # 兼容旧的单文件接口
        if "file" in request.files:
            uploaded_files = [request.files["file"]]
        else:
            return jsonify({"error": "没有导入文件"}), 400

    allowed_ext = {".png", ".jpg", ".jpeg", ".pdf"}
    task_id = str(uuid.uuid4())[:8]
    task_dir = UPLOAD_DIR / task_id
    task_dir.mkdir(exist_ok=True)

    saved_paths = []
    for f in uploaded_files:
        if f.filename == "":
            continue
        ext = Path(f.filename).suffix.lower()
        if ext not in allowed_ext:
            return jsonify({"error": f"不支持的文件类型: {ext}"}), 400
        save_path = task_dir / f"{len(saved_paths):03d}_{f.filename}"
        f.save(str(save_path))
        saved_paths.append(str(save_path))

    if not saved_paths:
        return jsonify({"error": "没有有效文件"}), 400

    # 是否启用元素分层
    decompose = request.form.get("decompose", "false").lower() == "true"
    decompose_mode = request.form.get("decompose_mode", "sam" if decompose else "none")
    if decompose_mode not in {"none", "sam", "qwen"}:
        return jsonify({"error": f"不支持的分层模式: {decompose_mode}"}), 400
    inpaint_backend = request.form.get("inpaint_backend", "lama").strip().lower()
    if inpaint_backend not in INPAINT_BACKENDS:
        return jsonify({"error": f"不支持的底图清理算法: {inpaint_backend}"}), 400
    if decompose and decompose_mode == "none":
        decompose_mode = "sam"
    decompose = decompose_mode != "none"
    try:
        qwen_num_layers = int(request.form.get("qwen_num_layers", "4"))
    except ValueError:
        return jsonify({"error": "qwen_num_layers 必须是数字"}), 400
    qwen_num_layers = max(3, min(8, qwen_num_layers))
    qwen_api_key = request.form.get("qwen_api_key", "").strip()
    output_path = allocate_output_path(saved_paths, decompose_mode)
    # 记录任务
    tasks[task_id] = {
        "status": "processing",
        "total_files": len(saved_paths),
        "total_slides": 0,
        "current_slide": 0,
        "current_step": "",
        "steps_done": [],  # 已完成步骤的历史
        "message": "导入完成，开始处理...",
        "decompose": decompose,
        "decompose_mode": decompose_mode,
        "inpaint_backend": inpaint_backend,
        "inpaint_backend_label": INPAINT_BACKEND_LABELS[inpaint_backend],
        "qwen_num_layers": qwen_num_layers,
        "qwen_api_configured": bool(qwen_api_key or os.environ.get("FAL_KEY")),
        "output": output_path,
        "error": None,
        "created_at": time.time(),
        "elapsed": 0,
    }

    # 在后台线程中处理
    thread = threading.Thread(
        target=_process_task,
        args=(
            task_id,
            saved_paths,
            str(task_dir),
            decompose,
            decompose_mode,
            qwen_num_layers,
            inpaint_backend,
            qwen_api_key,
        ),
        daemon=True,
    )
    thread.start()

    return jsonify({
        "task_id": task_id,
        "status": "processing",
        "decompose": decompose,
        "decompose_mode": decompose_mode,
        "inpaint_backend": inpaint_backend,
        "inpaint_backend_label": INPAINT_BACKEND_LABELS[inpaint_backend],
        "qwen_num_layers": qwen_num_layers,
        "qwen_api_configured": bool(qwen_api_key or os.environ.get("FAL_KEY")),
    })


@app.route("/api/download/<task_id>")
def download(task_id):
    """下载转换结果"""
    task = tasks.get(task_id)
    if not task:
        return jsonify({"error": "任务不存在"}), 404
    if task["status"] != "done":
        return jsonify({"error": "任务未完成"}), 400

    return send_file(
        task["output"],
        as_attachment=True,
        download_name=Path(task["output"]).name,
    )


@app.route("/api/status/<task_id>")
def status(task_id):
    """查询任务状态（前端轮询）"""
    task = tasks.get(task_id)
    if not task:
        return jsonify({"error": "任务不存在"}), 404
    if task.get("status") in {"processing", "generating"}:
        task["elapsed"] = round(time.time() - task.get("created_at", time.time()), 1)

    response = {
        "task_id": task_id,
        "status": task["status"],
        "message": task["message"],
        "current_slide": task.get("current_slide", 0),
        "total_slides": task.get("total_slides", 0),
        "current_step": task.get("current_step", ""),
        "decompose": task.get("decompose", False),
        "decompose_mode": task.get("decompose_mode", "none"),
        "inpaint_backend": task.get("inpaint_backend", "lama"),
        "inpaint_backend_label": task.get("inpaint_backend_label", "LaMa"),
        "qwen_num_layers": task.get("qwen_num_layers", 4),
        "qwen_api_configured": task.get("qwen_api_configured", False),
        "elapsed": task.get("elapsed", 0),
    }

    if task["status"] == "done":
        response["download_url"] = f"/api/download/{task_id}"
        if task.get("output"):
            response["output_name"] = Path(task["output"]).name
    elif task["status"] == "preview":
        response["preview_url"] = f"/api/preview/{task_id}"
    elif task["status"] == "error":
        response["error"] = task.get("error", "未知错误")

    return jsonify(response)


# ─── FastSAM 交互式预览接口 ───

@app.route("/api/preview/<task_id>")
def preview(task_id):
    """
    返回 FastSAM 分割预览数据。
    包含：可视化图 URL + 每个 mask 的元数据（编号、面积、bbox、轮廓点）
    """
    task = tasks.get(task_id)
    if not task:
        return jsonify({"error": "任务不存在"}), 404
    if task["status"] != "preview":
        return jsonify({"error": "任务不在预览状态"}), 400

    preview_data = task.get("preview_data", {})
    return jsonify(preview_data)


@app.route("/api/preview_image/<task_id>/<filename>")
def preview_image(task_id, filename):
    """提供预览相关的图片文件"""
    task = tasks.get(task_id)
    if not task:
        return jsonify({"error": "任务不存在"}), 404

    preview_dir = task.get("preview_dir")
    if not preview_dir:
        return jsonify({"error": "无预览数据"}), 404

    file_path = os.path.join(preview_dir, filename)
    if not os.path.exists(file_path):
        return jsonify({"error": "文件不存在"}), 404

    return send_file(file_path)


@app.route("/api/confirm/<task_id>", methods=["POST"])
def confirm(task_id):
    """
    用户确认分割结果，提交合并/删除决策，生成最终 PPTX。
    
    请求体 JSON:
    {
        "slides": [
            {
                "slide_index": 0,
                "keep": [0, 1, 4, 5],       // 保留的 mask 编号
                "merge": [[2, 3], [6, 7]],   // 需要合并的编号组
                "delete": [8, 9]             // 删除的编号
            },
            ...
        ]
    }
    """
    task = tasks.get(task_id)
    if not task:
        return jsonify({"error": "任务不存在"}), 404
    if task["status"] != "preview":
        return jsonify({"error": "任务不在预览状态"}), 400

    data = request.get_json()
    if not data:
        return jsonify({"error": "缺少请求体"}), 400

    # 保存用户决策
    task["user_decision"] = data
    task["status"] = "generating"
    task["message"] = "正在生成 PPTX..."

    # 后台生成 PPTX
    thread = threading.Thread(
        target=_generate_pptx_with_decision,
        args=(task_id,),
        daemon=True,
    )
    thread.start()

    return jsonify({"status": "generating", "message": "正在生成 PPTX..."})


def public_step_message(message: str) -> str:
    """把内部处理提示转换成面向用户的产品步骤文案。"""
    text = message or ""
    if "OCR" in text or "检测文字" in text:
        return "正在读取页面内容..."
    if "字体" in text:
        return "正在贴近原页面样式..."
    if "LaMa" in text or "本地均值" in text or "清理底图" in text or "擦除" in text:
        return "正在清理页面底图..."
    if "SAM" in text or "LayerD" in text or "Qwen" in text or "fal" in text or "分割" in text or "分层" in text:
        return "正在拆分页面元素..."
    if "生成" in text and "PPTX" in text:
        return "正在生成 PPTX..."
    return text


def _process_task(
    task_id: str,
    saved_paths: list,
    task_dir: str,
    decompose: bool,
    decompose_mode: str = "sam",
    qwen_num_layers: int = 4,
    inpaint_backend: str = "lama",
    qwen_api_key: str = "",
):
    """后台处理任务"""
    try:
        start_time = time.time()

        # 展开所有图片（PDF 拆页）
        tasks[task_id]["current_step"] = "解析文件..."
        all_images = []
        for fpath in saved_paths:
            if fpath.lower().endswith(".pdf"):
                pages = pdf_to_images(fpath, dpi=300, output_dir=task_dir)
                all_images.extend(pages)
            else:
                all_images.append(fpath)

        tasks[task_id]["total_slides"] = len(all_images)
        tasks[task_id]["message"] = f"共 {len(all_images)} 页待处理"

        # Qwen 云端分层直接生成 PPTX；FastSAM 分层保留交互式预览校对流程。
        if decompose_mode == "qwen":
            tasks[task_id]["message"] += f"（AI 智能分层，{qwen_num_layers} 层）"
            output_path = tasks[task_id]["output"]

            def progress_cb(current_idx, total, filename, msg):
                public_msg = public_step_message(msg)
                tasks[task_id]["current_slide"] = current_idx + 1
                tasks[task_id]["current_step"] = public_msg
                tasks[task_id]["message"] = f"[{current_idx+1}/{total}] {public_msg}"
                tasks[task_id]["elapsed"] = round(time.time() - start_time, 1)

            images_to_pptx(
                image_paths=all_images,
                output_path=output_path,
                device=DEFAULT_DEVICE,
                progress_callback=progress_cb,
                decompose=True,
                decompose_mode="qwen",
                qwen_num_layers=qwen_num_layers,
                qwen_api_key=qwen_api_key,
                inpaint_backend=inpaint_backend,
            )

            tasks[task_id]["status"] = "done"
            tasks[task_id]["output"] = output_path
            tasks[task_id]["elapsed"] = round(time.time() - start_time, 1)
            tasks[task_id]["message"] = f"转换完成，耗时 {tasks[task_id]['elapsed']:.0f}s"
            compact_finished_task(tasks[task_id])

        # 如果启用 FastSAM 分层模式，走预览流程
        elif decompose:
            tasks[task_id]["message"] += "（分层模式）"
            _process_for_preview(task_id, all_images, task_dir, start_time, inpaint_backend=inpaint_backend)
        else:
            # 标准模式：直接生成 PPTX
            output_path = tasks[task_id]["output"]

            def progress_cb(current_idx, total, filename, msg):
                public_msg = public_step_message(msg)
                tasks[task_id]["current_slide"] = current_idx + 1
                tasks[task_id]["current_step"] = public_msg
                tasks[task_id]["message"] = f"[{current_idx+1}/{total}] {public_msg}"
                tasks[task_id]["elapsed"] = round(time.time() - start_time, 1)

            images_to_pptx(
                image_paths=all_images,
                output_path=output_path,
                device=DEFAULT_DEVICE,
                progress_callback=progress_cb,
                decompose=False,
                inpaint_backend=inpaint_backend,
            )

            tasks[task_id]["status"] = "done"
            tasks[task_id]["output"] = output_path
            tasks[task_id]["elapsed"] = round(time.time() - start_time, 1)
            tasks[task_id]["message"] = f"转换完成，耗时 {tasks[task_id]['elapsed']:.0f}s"
            compact_finished_task(tasks[task_id])

    except Exception as e:
        tasks[task_id]["status"] = "error"
        tasks[task_id]["error"] = str(e)
        tasks[task_id]["message"] = f"处理失败: {e}"
    finally:
        release_memory()


def _process_for_preview(task_id: str, all_images: list, task_dir: str, start_time: float, inpaint_backend: str = "lama"):
    """
    分层模式：先生成页面底图和元素候选，再进入预览状态等待用户确认。
    """
    import colorsys

    preview_dir = os.path.join(task_dir, "preview")
    os.makedirs(preview_dir, exist_ok=True)

    slides_data = []

    for idx, img_path in enumerate(all_images):
        img_path = constrain_image_for_processing(img_path, output_dir=preview_dir)
        tasks[task_id]["current_slide"] = idx + 1
        tasks[task_id]["elapsed"] = round(time.time() - start_time, 1)

        tasks[task_id]["current_step"] = "正在读取页面内容..."
        tasks[task_id]["message"] = f"[{idx+1}/{len(all_images)}] 正在读取页面内容..."
        blocks = detect_text_paddle(img_path)
        if blocks:
            blocks = match_fonts_for_blocks(img_path, blocks)

        tasks[task_id]["current_step"] = "正在清理页面底图..."
        tasks[task_id]["message"] = f"[{idx+1}/{len(all_images)}] 正在清理页面底图..."
        if blocks:
            background = remove_text_from_image(img_path, blocks, device=DEFAULT_DEVICE, backend=inpaint_backend)
        else:
            background = Image.open(img_path).convert("RGB")

        # 保存干净背景图
        bg_clean_path = os.path.join(preview_dir, f"slide_{idx:03d}_bg.png")
        background.save(bg_clean_path)

        tasks[task_id]["current_step"] = "正在拆分页面元素..."
        tasks[task_id]["message"] = f"[{idx+1}/{len(all_images)}] 正在拆分页面元素..."

        bg_rgb = background.convert("RGB")
        bg_np = np.array(bg_rgb)
        img_w, img_h = bg_rgb.size
        total_area = img_w * img_h

        try:
            masks = generate_background_fastsam_masks(bg_np, device=DEFAULT_DEVICE)
        except Exception as e:
            print(f"  [FastSAM] 预览分割失败，改用 OpenCV 兜底分割: {e}", flush=True)
            masks = segment_background_cv_masks(bg_np)

        # 后处理（与 decompose_background_sam 一致）
        min_area_ratio = 0.001
        max_area_ratio = 0.80
        bg_color_threshold = 30.0
        overlap_iou_threshold = 0.7

        filtered = [m for m in masks if min_area_ratio <= m["area"] / total_area <= max_area_ratio]

        corner_size = max(20, min(img_w, img_h) // 10)
        corners = np.concatenate([
            bg_np[:corner_size, :corner_size].reshape(-1, 3),
            bg_np[:corner_size, -corner_size:].reshape(-1, 3),
            bg_np[-corner_size:, :corner_size].reshape(-1, 3),
            bg_np[-corner_size:, -corner_size:].reshape(-1, 3),
        ], axis=0)
        bg_color = np.median(corners, axis=0)

        color_filtered = []
        for m in filtered:
            seg = m["segmentation"]
            mask_mean_color = np.mean(bg_np[seg], axis=0)
            color_diff = np.sqrt(np.sum((mask_mean_color - bg_color) ** 2))
            if color_diff >= bg_color_threshold:
                color_filtered.append(m)

        color_filtered.sort(key=lambda x: x["area"], reverse=True)
        keep = []
        for m in color_filtered:
            should_keep = True
            seg_m = m["segmentation"]
            for kept in keep:
                seg_k = kept["segmentation"]
                intersection = np.logical_and(seg_m, seg_k).sum()
                union = np.logical_or(seg_m, seg_k).sum()
                if union > 0 and intersection / union > overlap_iou_threshold:
                    should_keep = False
                    break
            if should_keep:
                keep.append(m)

        # 包含关系过滤：如果大 mask 被其内部所有小 mask 合计覆盖了 90%+ 面积，才删掉大的
        # （说明大 mask 只是小 mask 们的"容器"，没有独立内容）
        containment_filtered = []
        for i, m in enumerate(keep):
            seg_m = m["segmentation"]
            area_m = m["area"]
            # 找出所有被 m 包含的小 mask（小 mask 的 90%+ 在 m 内部）
            total_covered = np.zeros_like(seg_m, dtype=bool)
            has_children = False
            for j, other in enumerate(keep):
                if i == j:
                    continue
                if other["area"] >= area_m:
                    continue  # 只看比自己小的
                seg_other = other["segmentation"]
                overlap = np.logical_and(seg_m, seg_other).sum()
                if other["area"] > 0 and overlap / other["area"] > 0.9:
                    # other 是 m 的"子元素"
                    total_covered = np.logical_or(total_covered, seg_other)
                    has_children = True
            if has_children:
                # 计算所有子 mask 合计覆盖了大 mask 多少面积
                coverage = total_covered.sum() / area_m if area_m > 0 else 0
                if coverage > 0.8:
                    # 大 mask 80%+ 被子 mask 们覆盖了，它只是容器，删掉
                    continue
            containment_filtered.append(m)
        keep = containment_filtered

        # 生成可视化图
        tasks[task_id]["current_step"] = "生成预览..."
        tasks[task_id]["message"] = f"[{idx+1}/{len(all_images)}] 生成预览..."

        vis = bg_rgb.copy().convert("RGBA")
        overlay = Image.new("RGBA", (img_w, img_h), (0, 0, 0, 0))

        for midx, m in enumerate(keep):
            hue = midx / max(len(keep), 1)
            r, g, b = colorsys.hsv_to_rgb(hue, 0.7, 0.9)
            color = (int(r * 255), int(g * 255), int(b * 255), 80)
            seg = m["segmentation"]
            mask_img = Image.fromarray((seg * 255).astype(np.uint8), mode="L")
            colored = Image.new("RGBA", (img_w, img_h), color)
            overlay = Image.composite(colored, overlay, mask_img)

        vis = Image.alpha_composite(vis, overlay).convert("RGB")
        viz_path = os.path.join(preview_dir, f"slide_{idx:03d}_viz.jpg")
        vis.save(viz_path, "JPEG", quality=90)

        # 构造 mask 元数据（不含 segmentation 原始数据，太大）
        # 保存 segmentation 到 npz 文件供后续生成 PPTX 用
        masks_npz_path = os.path.join(preview_dir, f"slide_{idx:03d}_masks.npz")
        seg_arrays = {f"mask_{midx}": m["segmentation"].astype(np.uint8) for midx, m in enumerate(keep)}
        np.savez_compressed(masks_npz_path, **seg_arrays)

        # mask 元数据（发给前端）— 包含 polygon 点集供 Fabric.js 渲染
        mask_meta = []
        for midx, m in enumerate(keep):
            hue = midx / max(len(keep), 1)
            r, g, b = colorsys.hsv_to_rgb(hue, 0.7, 0.9)

            # 将 mask 转为 polygon 点集（用于前端 Canvas 精确渲染和 hit-test）
            seg_uint8 = (m["segmentation"] * 255).astype(np.uint8)
            contours, _ = cv2.findContours(seg_uint8, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
            polygons = []
            for contour in contours:
                if len(contour) < 3:
                    continue
                # 简化轮廓减少数据量
                epsilon = 0.002 * cv2.arcLength(contour, True)
                approx = cv2.approxPolyDP(contour, epsilon, True)
                points = approx.reshape(-1, 2).tolist()
                if len(points) >= 3:
                    polygons.append(points)

            mask_meta.append({
                "id": midx,
                "area": m["area"],
                "area_pct": round(m["area"] / total_area * 100, 1),
                "bbox": m["bbox"],  # [x, y, w, h]
                "color": [int(r * 255), int(g * 255), int(b * 255)],
                "polygons": polygons,
            })

        slides_data.append({
            "slide_index": idx,
            "image_path": img_path,
            "bg_clean": f"slide_{idx:03d}_bg.png",
            "viz_image": f"slide_{idx:03d}_viz.jpg",
            "masks_file": f"slide_{idx:03d}_masks.npz",
            "img_width": img_w,
            "img_height": img_h,
            "masks": mask_meta,
            "blocks_count": len(blocks),
        })

        # 保存 blocks 数据供后续 PPTX 生成
        import pickle
        blocks_path = os.path.join(preview_dir, f"slide_{idx:03d}_blocks.pkl")
        with open(blocks_path, "wb") as f:
            pickle.dump(blocks, f)

        del masks, filtered, color_filtered, keep, seg_arrays, bg_np, bg_rgb, background
        release_memory()

    # 进入预览状态
    tasks[task_id]["status"] = "preview"
    tasks[task_id]["preview_dir"] = preview_dir
    tasks[task_id]["preview_data"] = {
        "task_id": task_id,
        "total_slides": len(all_images),
        "slides": slides_data,
    }
    tasks[task_id]["all_images"] = all_images
    tasks[task_id]["elapsed"] = round(time.time() - start_time, 1)
    tasks[task_id]["message"] = f"分割完成，请确认元素分组"
    tasks[task_id]["current_step"] = "等待确认..."


def _generate_pptx_with_decision(task_id: str):
    """根据用户的合并/删除决策生成最终 PPTX"""
    import pickle
    from engine import create_deduped_rgba_layers, create_pptx_with_layers, PIXELS_TO_EMU
    from pptx import Presentation
    from pptx.util import Emu

    try:
        task = tasks[task_id]
        start_time = time.time()
        preview_dir = task["preview_dir"]
        preview_data = task["preview_data"]
        decision = task["user_decision"]
        all_images = task["all_images"]

        output_path = task["output"]
        prs = Presentation()

        slides_decision = decision.get("slides", [])

        for slide_info in preview_data["slides"]:
            idx = slide_info["slide_index"]
            img_w = slide_info["img_width"]
            img_h = slide_info["img_height"]

            # 设置幻灯片尺寸
            slide_width = img_w * PIXELS_TO_EMU
            slide_height = img_h * PIXELS_TO_EMU
            if idx == 0:
                prs.slide_width = Emu(slide_width)
                prs.slide_height = Emu(slide_height)
            else:
                if slide_width > prs.slide_width:
                    prs.slide_width = Emu(slide_width)
                if slide_height > prs.slide_height:
                    prs.slide_height = Emu(slide_height)

            # 加载 masks
            masks_npz_path = os.path.join(preview_dir, slide_info["masks_file"])
            with np.load(masks_npz_path) as masks_data:
                # 加载 blocks
                blocks_path = os.path.join(preview_dir, f"slide_{idx:03d}_blocks.pkl")
                with open(blocks_path, "rb") as f:
                    blocks = pickle.load(f)

                # 加载干净背景图
                bg_path = os.path.join(preview_dir, slide_info["bg_clean"])
                bg_image = Image.open(bg_path).convert("RGB")

                # 获取该 slide 的用户决策
                slide_dec = None
                for sd in slides_decision:
                    if sd.get("slide_index") == idx:
                        slide_dec = sd
                        break

                # 根据决策构建最终 mask 列表
                if slide_dec:
                    keep_ids = set(slide_dec.get("keep", []))
                    merge_groups = slide_dec.get("merge", [])
                    delete_ids = set(slide_dec.get("delete", []))
                else:
                    # 没有决策 = 保留所有
                    total_masks = len([k for k in masks_data.files])
                    keep_ids = set(range(total_masks))
                    merge_groups = []
                    delete_ids = set()

                # 构建最终的 segmentation 列表
                final_segs = []

                # 先处理合并组
                merged_ids = set()
                for group in merge_groups:
                    merged_seg = np.zeros((img_h, img_w), dtype=np.uint8)
                    for mid in group:
                        key = f"mask_{mid}"
                        if key in masks_data:
                            merged_seg = np.logical_or(merged_seg, masks_data[key]).astype(np.uint8)
                            merged_ids.add(mid)
                    if merged_seg.any():
                        final_segs.append(merged_seg)

                # 再处理独立保留的
                for mid in keep_ids:
                    if mid in merged_ids or mid in delete_ids:
                        continue
                    key = f"mask_{mid}"
                    if key in masks_data:
                        final_segs.append(masks_data[key].copy())

            # 生成 RGBA 图层列表
            bg_np = np.array(bg_image)
            layers = create_deduped_rgba_layers(bg_np, final_segs)

            # 生成 slide
            import tempfile as _tempfile
            bg_dir = _tempfile.mkdtemp(prefix="decklens_confirm_")

            try:
                layer_info = create_pptx_with_layers(
                    layers, blocks, "", bg_dir=bg_dir, slide_index=idx
                )

                blank_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_layout)

                for path, x, y, w, h in layer_info:
                    slide.shapes.add_picture(
                        path,
                        Emu(x * PIXELS_TO_EMU),
                        Emu(y * PIXELS_TO_EMU),
                        width=Emu(w * PIXELS_TO_EMU),
                        height=Emu(h * PIXELS_TO_EMU),
                    )

                # 文字框
                from pptx.util import Pt
                from pptx.dml.color import RGBColor
                from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
                from pptx.oxml.ns import qn
                from lxml import etree as _etree

                for block in blocks:
                    left = block.x * PIXELS_TO_EMU
                    top = block.y * PIXELS_TO_EMU
                    width = block.w * PIXELS_TO_EMU
                    height = block.h * PIXELS_TO_EMU

                    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
                    tf = txBox.text_frame
                    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    tf.word_wrap = False
                    tf.margin_left = Emu(0)
                    tf.margin_right = Emu(0)
                    tf.margin_top = Emu(0)
                    tf.margin_bottom = Emu(0)

                    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
                    body_pr = tf._txBody.bodyPr
                    body_pr.set('anchor', 'ctr')

                    p = tf.paragraphs[0]
                    p.text = block.text

                    import re as _re
                    pPr = p._p.get_or_add_pPr()
                    if _re.match(r'^[0-9A-Za-z\s\-_\.]+$', block.text.strip()):
                        line_spacing_val = '85000'
                    else:
                        line_spacing_val = '100000'
                    lnSpc = _etree.SubElement(pPr, qn('a:lnSpc'))
                    spcPct = _etree.SubElement(lnSpc, qn('a:spcPct'))
                    spcPct.set('val', line_spacing_val)
                    spcBef = _etree.SubElement(pPr, qn('a:spcBef'))
                    spcPts = _etree.SubElement(spcBef, qn('a:spcPts'))
                    spcPts.set('val', '0')
                    spcAft = _etree.SubElement(pPr, qn('a:spcAft'))
                    spcPts2 = _etree.SubElement(spcAft, qn('a:spcPts'))
                    spcPts2.set('val', '0')

                    run = p.runs[0]
                    run.font.size = Pt(int(block.font_size_pt))
                    r, g, b = block.color
                    run.font.color.rgb = RGBColor(r, g, b)
                    run.font.bold = block.bold
                    run.font.name = block.font_name
                    try:
                        rPr = run._r.get_or_add_rPr()
                        ea = _etree.SubElement(rPr, qn('a:ea'))
                        ea.set('typeface', block.font_name)
                    except Exception:
                        pass

                    txBox.fill.background()
                    txBox.line.fill.background()
            finally:
                shutil.rmtree(bg_dir, ignore_errors=True)

            del masks_data, bg_image, bg_np, layers, final_segs
            release_memory()

        prs.save(output_path)

        task["status"] = "done"
        task["output"] = output_path
        task["elapsed"] = round(task.get("elapsed", 0) + (time.time() - start_time), 1)
        task["message"] = f"生成完成，耗时 {task['elapsed']:.0f}s"
        compact_finished_task(task)

    except Exception as e:
        tasks[task_id]["status"] = "error"
        tasks[task_id]["error"] = str(e)
        tasks[task_id]["message"] = f"生成失败: {e}"
        import traceback
        traceback.print_exc()
    finally:
        release_memory()


if __name__ == "__main__":
    port = int(os.environ.get("PORT", "8080"))
    print("=" * 50)
    print("  DeckLens — 图片化页面转可编辑 PPTX 与元素分层工作台")
    print("  支持多文件批量处理 | 元素分层 | 实时进度显示")
    print(f"  访问: http://localhost:{port}")
    print("=" * 50)
    app.run(host="0.0.0.0", port=port, debug=False)
